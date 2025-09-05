import os
import sys
import logging
import signal
import traceback
import time
import argparse
import uuid
import asyncio
from datetime import datetime
from fastapi import FastAPI, HTTPException, Request, File, UploadFile, Form, Depends, Query, BackgroundTasks
from fastapi.responses import JSONResponse, Response, HTMLResponse
from fastapi.middleware.cors import CORSMiddleware
import tempfile
import shutil
from pathlib import Path
from typing import List, Optional, Dict, Any
import json
from pydantic import BaseModel, Field
from io import BytesIO
import PyPDF2
import xml.etree.ElementTree as ET

# Import Speech2Text
from Speech2Text import Speech2Text

# Import Text2Speech
from Text2Speech import Text2Speech

# LightRAG imports
try:
    from lightrag import LightRAG, QueryParam
    from lightrag.llm.openai import gpt_4o_mini_complete, openai_embed, openai_complete_if_cache, gpt_4o_complete, openai_complete
    from lightrag.llm.ollama import ollama_model_complete, ollama_embed
    from lightrag.utils import EmbeddingFunc, setup_logger
    from lightrag.kg.shared_storage import initialize_pipeline_status

    LIGHTRAG_AVAILABLE = True
except ImportError as e:
    logging.warning(f"LightRAG not available: {e}")
    LIGHTRAG_AVAILABLE = False

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[logging.StreamHandler()]
)
logger = logging.getLogger("clara-backend")

# Store start time
START_TIME = datetime.now().isoformat()

# Parse command line arguments
parser = argparse.ArgumentParser(description='Clara Backend Server')
parser.add_argument('--host', type=str, default='127.0.0.1', help='Host to bind to')
parser.add_argument('--port', type=int, default=5000, help='Port to bind to')
args = parser.parse_args()

# Use the provided host and port
HOST = args.host
PORT = args.port

logger.info(f"Starting server on {HOST}:{PORT}")

# Setup FastAPI
app = FastAPI(title="Clara Backend API", version="1.0.0")

# Import and include the diffusers API router
# Add CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # In production, replace with specific origins
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
    expose_headers=["*"]
)

# Add global exception middleware
@app.middleware("http")
async def catch_exceptions_middleware(request: Request, call_next):
    try:
        return await call_next(request)
    except Exception as e:
        logger.error(f"Request to {request.url} failed: {str(e)}")
        logger.error(traceback.format_exc())
        return JSONResponse(
            status_code=500,
            content={"error": str(e), "detail": traceback.format_exc()}
        )

# Database path in user's home directory for persistence
home_dir = os.path.expanduser("~")
data_dir = os.path.join(home_dir, ".clara")
os.makedirs(data_dir, exist_ok=True)

# LightRAG Configuration and Storage
if LIGHTRAG_AVAILABLE:
    # Environment setup for LightRAG
    os.environ["NEO4J_URI"] = os.getenv("NEO4J_URI", "neo4j://localhost:7687")
    os.environ["NEO4J_USERNAME"] = os.getenv("NEO4J_USERNAME", "neo4j")
    os.environ["NEO4J_PASSWORD"] = os.getenv("NEO4J_PASSWORD", "password")
    os.environ["OPENAI_API_KEY"] = os.getenv("OPENAI_API_KEY", "your-openai-api-key")

    # Setup LightRAG logging
    setup_logger("lightrag", level="INFO")
    
    # Storage paths for LightRAG
    LIGHTRAG_STORAGE_PATH = Path(os.path.join(data_dir, "lightrag_storage"))
    LIGHTRAG_STORAGE_PATH.mkdir(exist_ok=True)
    LIGHTRAG_METADATA_PATH = LIGHTRAG_STORAGE_PATH / "metadata"
    LIGHTRAG_METADATA_PATH.mkdir(exist_ok=True)

    # Global storage for LightRAG notebooks and documents
    lightrag_notebooks_db: Dict[str, Dict] = {}
    lightrag_documents_db: Dict[str, Dict] = {}
    lightrag_instances: Dict[str, LightRAG] = {}
    # Chat history storage for maintaining conversation context
    chat_history_db: Dict[str, List[Dict]] = {}  # notebook_id -> [messages]

    # Persistence files
    NOTEBOOKS_DB_FILE = LIGHTRAG_METADATA_PATH / "notebooks.json"
    DOCUMENTS_DB_FILE = LIGHTRAG_METADATA_PATH / "documents.json"
    CHAT_HISTORY_DB_FILE = LIGHTRAG_METADATA_PATH / "chat_history.json"

    def save_notebooks_db():
        """Save notebooks database to disk"""
        try:
            # Convert datetime objects to ISO strings for JSON serialization
            serializable_data = {}
            for notebook_id, notebook_data in lightrag_notebooks_db.items():
                serializable_notebook = notebook_data.copy()
                if isinstance(serializable_notebook.get('created_at'), datetime):
                    serializable_notebook['created_at'] = serializable_notebook['created_at'].isoformat()
                serializable_data[notebook_id] = serializable_notebook
            
            with open(NOTEBOOKS_DB_FILE, 'w') as f:
                json.dump(serializable_data, f, indent=2)
            logger.info(f"Saved {len(serializable_data)} notebooks to {NOTEBOOKS_DB_FILE}")
        except Exception as e:
            logger.error(f"Error saving notebooks database: {e}")

    def load_notebooks_db():
        """Load notebooks database from disk"""
        global lightrag_notebooks_db
        try:
            if NOTEBOOKS_DB_FILE.exists():
                with open(NOTEBOOKS_DB_FILE, 'r') as f:
                    data = json.load(f)
                
                # Convert ISO strings back to datetime objects
                for notebook_id, notebook_data in data.items():
                    if isinstance(notebook_data.get('created_at'), str):
                        notebook_data['created_at'] = datetime.fromisoformat(notebook_data['created_at'])
                
                lightrag_notebooks_db = data
                logger.info(f"Loaded {len(data)} notebooks from {NOTEBOOKS_DB_FILE}")
            else:
                logger.info("No existing notebooks database found")
        except Exception as e:
            logger.error(f"Error loading notebooks database: {e}")
            lightrag_notebooks_db = {}

    def save_documents_db():
        """Save documents database to disk"""
        try:
            # Convert datetime objects to ISO strings for JSON serialization
            serializable_data = {}
            for document_id, document_data in lightrag_documents_db.items():
                serializable_document = document_data.copy()
                # Convert all datetime values to ISO strings
                for key, value in serializable_document.items():
                    if isinstance(value, datetime):
                        serializable_document[key] = value.isoformat()
                serializable_data[document_id] = serializable_document
            
            with open(DOCUMENTS_DB_FILE, 'w') as f:
                json.dump(serializable_data, f, indent=2)
            logger.info(f"Saved {len(serializable_data)} documents to {DOCUMENTS_DB_FILE}")
        except Exception as e:
            logger.error(f"Error saving documents database: {e}")

    def load_documents_db():
        """Load documents database from disk"""
        global lightrag_documents_db
        try:
            if DOCUMENTS_DB_FILE.exists():
                with open(DOCUMENTS_DB_FILE, 'r') as f:
                    data = json.load(f)
                
                # Convert ISO strings back to datetime objects
                for document_id, document_data in data.items():
                    for key, value in document_data.items():
                        if isinstance(value, str) and key.endswith('_at'):
                            try:
                                document_data[key] = datetime.fromisoformat(value)
                            except ValueError:
                                pass  # Keep as string if not a valid ISO datetime
                
                lightrag_documents_db = data
                logger.info(f"Loaded {len(data)} documents from {DOCUMENTS_DB_FILE}")
            else:
                logger.info("No existing documents database found")
        except Exception as e:
            logger.error(f"Error loading documents database: {e}")
            lightrag_documents_db = {}

    def save_chat_history_db():
        """Save chat history database to disk"""
        try:
            # Convert datetime objects to ISO strings for JSON serialization
            serializable_data = {}
            for notebook_id, messages in chat_history_db.items():
                serializable_messages = []
                for message in messages:
                    serializable_message = message.copy()
                    if isinstance(serializable_message.get('timestamp'), datetime):
                        serializable_message['timestamp'] = serializable_message['timestamp'].isoformat()
                    serializable_messages.append(serializable_message)
                serializable_data[notebook_id] = serializable_messages
            
            with open(CHAT_HISTORY_DB_FILE, 'w') as f:
                json.dump(serializable_data, f, indent=2)
            logger.info(f"Saved chat history for {len(serializable_data)} notebooks to {CHAT_HISTORY_DB_FILE}")
        except Exception as e:
            logger.error(f"Error saving chat history database: {e}")

    def load_chat_history_db():
        """Load chat history database from disk"""
        global chat_history_db
        try:
            if CHAT_HISTORY_DB_FILE.exists():
                with open(CHAT_HISTORY_DB_FILE, 'r') as f:
                    data = json.load(f)
                
                # Convert ISO strings back to datetime objects
                for notebook_id, messages in data.items():
                    for message in messages:
                        if isinstance(message.get('timestamp'), str):
                            try:
                                message['timestamp'] = datetime.fromisoformat(message['timestamp'])
                            except ValueError:
                                pass  # Keep as string if not a valid ISO datetime
                
                chat_history_db = data
                logger.info(f"Loaded chat history for {len(data)} notebooks from {CHAT_HISTORY_DB_FILE}")
            else:
                logger.info("No existing chat history database found")
        except Exception as e:
            logger.error(f"Error loading chat history database: {e}")
            chat_history_db = {}

    # Load existing data on startup
    load_notebooks_db()
    load_documents_db()
    load_chat_history_db()

# Speech2Text instance cache
speech2text_instance = None

def get_speech2text():
    """Create or retrieve the Speech2Text instance from cache"""
    global speech2text_instance
    
    if speech2text_instance is None:
        # Use tiny model with CPU for maximum compatibility
        speech2text_instance = Speech2Text(
            model_size="tiny",
            device="cpu",
            compute_type="int8"
        )
    
    return speech2text_instance

# Text2Speech instance cache
text2speech_instance = None

def get_text2speech():
    """Create or retrieve the Text2Speech instance from cache"""
    global text2speech_instance
    
    if text2speech_instance is None:
        # Initialize with auto engine selection (will prefer Kokoro if available)
        text2speech_instance = Text2Speech(
            engine="auto",
            language="en",
            slow=False,
            voice="af_sarah",
            speed=1.0
        )
    
    return text2speech_instance

# LightRAG Utility Functions
if LIGHTRAG_AVAILABLE:
    def create_llm_func(provider_config: Dict[str, Any]):
        """Create LLM function based on provider configuration using LightRAG built-in functions"""
        provider_type = provider_config.get('type', 'ollama')
        model_name = provider_config.get('model', 'gemma3:4b')
        api_key = provider_config.get('apiKey', '')
        base_url = provider_config.get('baseUrl', '')
        
        # Validate API key for non-Ollama providers
        if provider_type != 'ollama' and (not api_key or api_key.strip() == 'your-api-key'):
            raise ValueError(f"Invalid or missing API key for provider type: {provider_type}")
        
        try:
            if provider_type == 'openai':
                logger.info(f"Using OpenAI LLM with model: {model_name}")
                if 'gpt-4o-mini' in model_name:
                    return gpt_4o_mini_complete
                elif 'gpt-4o' in model_name:
                    return gpt_4o_complete
                else:
                    return openai_complete
            
            elif provider_type == 'openai_compatible':
                logger.info(f"Using OpenAI-compatible LLM: {base_url} with model: {model_name}")
                # Use LightRAG's built-in function directly - it will handle the calling
                return openai_complete_if_cache
            
            elif provider_type == 'ollama':
                logger.info(f"Using Ollama LLM with model: {model_name}")
                # Use LightRAG's built-in function directly - it will handle the calling
                return ollama_model_complete
            
            else:
                raise ValueError(f"Unsupported provider type: {provider_type}")
                        
        except Exception as e:
            logger.error(f"Error creating LLM function for provider {provider_type}: {e}")
            raise

    def create_embedding_func(provider_config: Dict[str, Any]):
        """Create embedding function based on provider configuration using LightRAG built-in functions"""
        provider_type = provider_config.get('type', 'ollama')
        model_name = provider_config.get('model', 'mxbai-embed-large:latest')
        api_key = provider_config.get('apiKey', '')
        base_url = provider_config.get('baseUrl', '')
        
        # Validate API key for non-Ollama providers
        if provider_type != 'ollama' and (not api_key or api_key.strip() == 'your-api-key'):
            raise ValueError(f"Invalid or missing API key for embedding provider type: {provider_type}")
        
        try:
            if provider_type == 'openai':
                logger.info(f"Using OpenAI embeddings with model: {model_name}")
                # Use lambda function to pass the texts parameter correctly
                return lambda texts: openai_embed(
                    texts=texts,
                    model=model_name,
                    api_key=api_key.strip()
                )
            
            elif provider_type == 'openai_compatible':
                logger.info(f"Using OpenAI-compatible embeddings: {base_url} with model: {model_name}")
                # Try user's config first, fallback to Ollama
                def openai_compatible_embedding_func(texts):
                    try:
                        logger.info(f"Attempting embeddings with {base_url}")
                        result = openai_embed(
                            texts=texts,
                                model=model_name,
                            api_key=api_key.strip(),
                            base_url=base_url
                        )
                        logger.info(f"Successfully generated {len(result)} embeddings with OpenAI-compatible API")
                        return result
                    except Exception as e:
                        logger.warning(f"OpenAI-compatible embedding failed: {e}")
                        logger.info("Falling back to local Ollama embeddings")
                        try:
                            result = ollama_embed(
                                texts=texts,
                                embed_model="mxbai-embed-large:latest",
                                host="http://localhost:11434"
                            )
                            logger.info(f"Successfully generated {len(result)} embeddings with Ollama fallback")
                            return result
                        except Exception as fallback_error:
                            error_msg = f"Both OpenAI-compatible and Ollama embedding failed. OpenAI error: {e}, Ollama error: {fallback_error}"
                            logger.error(error_msg)
                            raise Exception(error_msg)
                
                return openai_compatible_embedding_func
            
            elif provider_type == 'ollama':
                logger.info(f"Using Ollama embeddings with model: {model_name}")
                # Use lambda function to pass parameters correctly
                return lambda texts: ollama_embed(
                    texts=texts,
                    embed_model=model_name,
                    host=base_url if base_url else "http://localhost:11434"
                )
            
            else:
                raise ValueError(f"Unsupported embedding provider type: {provider_type}")
                        
        except Exception as e:
            logger.error(f"Error creating embedding function for provider {provider_type}: {e}")
            raise

    def extract_text_from_pdf_lightrag(pdf_bytes: bytes) -> str:
        """Extract text from PDF bytes for LightRAG"""
        try:
            pdf_reader = PyPDF2.PdfReader(BytesIO(pdf_bytes))
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            return text
        except Exception as e:
            logger.error(f"Error extracting text from PDF: {e}")
            raise HTTPException(status_code=400, detail=f"Error processing PDF: {str(e)}")

    # Additional format support endpoint
    @app.get("/notebooks/supported-formats")
    async def get_supported_formats():
        """Get information about supported file formats for notebook documents"""
        basic_formats = [
            {"extension": "pdf", "description": "PDF documents", "library": "PyPDF2 (built-in)"},
            {"extension": "txt", "description": "Plain text files", "library": "built-in"},
            {"extension": "md", "description": "Markdown files", "library": "built-in"},
            {"extension": "csv", "description": "Comma-separated values", "library": "built-in"},
            {"extension": "json", "description": "JSON data files", "library": "built-in"},
            {"extension": "xml", "description": "XML documents", "library": "built-in"},
            {"extension": "html", "description": "HTML documents", "library": "built-in"},
            {"extension": "htm", "description": "HTML documents", "library": "built-in"}
        ]
        
        enhanced_formats = [
            {"extension": "docx", "description": "Word documents (newer)", "library": "python-docx", "install": "pip install python-docx"},
            {"extension": "doc", "description": "Word documents (legacy)", "library": "python-docx", "install": "pip install python-docx"},
            {"extension": "xlsx", "description": "Excel spreadsheets (newer)", "library": "pandas + openpyxl", "install": "pip install pandas openpyxl"},
            {"extension": "xls", "description": "Excel spreadsheets (legacy)", "library": "pandas + xlrd", "install": "pip install pandas xlrd"},
            {"extension": "pptx", "description": "PowerPoint presentations (newer)", "library": "python-pptx", "install": "pip install python-pptx"},
            {"extension": "ppt", "description": "PowerPoint presentations (legacy)", "library": "python-pptx", "install": "pip install python-pptx"},
            {"extension": "rtf", "description": "Rich Text Format", "library": "striprtf", "install": "pip install striprtf"}
        ]
        
        textract_formats = [
            {"extension": "epub", "description": "eBook format", "library": "textract", "install": "pip install textract"},
            {"extension": "odt", "description": "OpenDocument Text", "library": "textract", "install": "pip install textract"},
            {"extension": "ods", "description": "OpenDocument Spreadsheet", "library": "textract", "install": "pip install textract"},
            {"extension": "odp", "description": "OpenDocument Presentation", "library": "textract", "install": "pip install textract"}
        ]
        
        return {
            "basic_support": [f".{fmt['extension']}" for fmt in basic_formats],
            "enhanced_support": [f".{fmt['extension']}" for fmt in enhanced_formats],
            "textract_support": [f".{fmt['extension']}" for fmt in textract_formats],
            "details": {
                "basic": basic_formats,
                "enhanced": enhanced_formats,
                "textract": textract_formats
            },
            "installation_notes": {
                "basic": "These formats are supported out of the box",
                "enhanced": "Install additional libraries for enhanced document support",
                "textract": "Install textract for additional format support (requires system dependencies)"
            },
            "recommended_libraries": [
                "pip install python-docx python-pptx pandas openpyxl xlrd striprtf beautifulsoup4",
                "pip install textract  # Optional for additional formats"
            ]
        }

    async def extract_text_from_file(filename: str, file_content: bytes) -> str:
        """Extract text from various file formats supported by LightRAG"""
        try:
            file_ext = filename.lower().split('.')[-1] if '.' in filename else ''
            
            # PDF files
            if file_ext == 'pdf':
                return extract_text_from_pdf_lightrag(file_content)
            
            # Plain text files
            elif file_ext in ['txt', 'md', 'markdown', 'rst']:
                return file_content.decode('utf-8')
            
            # CSV files
            elif file_ext == 'csv':
                return file_content.decode('utf-8')
            
            # JSON files
            elif file_ext == 'json':
                import json
                json_data = json.loads(file_content.decode('utf-8'))
                return json.dumps(json_data, indent=2)
            
            # XML/HTML files
            elif file_ext in ['xml', 'html', 'htm']:
                content = file_content.decode('utf-8')
                try:
                    from bs4 import BeautifulSoup
                    if file_ext in ['html', 'htm']:
                        soup = BeautifulSoup(content, 'html.parser')
                        return soup.get_text(separator='\n', strip=True)
                    else:
                        soup = BeautifulSoup(content, 'xml')
                        return soup.get_text(separator='\n', strip=True)
                except ImportError:
                    # Fallback to raw content if BeautifulSoup not available
                    return content
            
            # Word documents
            elif file_ext in ['docx', 'doc']:
                try:
                    import docx
                    from io import BytesIO
                    doc = docx.Document(BytesIO(file_content))
                    content_parts = []
                    for para in doc.paragraphs:
                        text = para.text.strip()
                        if text:
                            content_parts.append(text)
                    return "\n\n".join(content_parts)
                except ImportError:
                    logger.warning("python-docx not available, cannot process Word documents")
                    raise HTTPException(
                        status_code=400, 
                        detail="Word document processing requires python-docx library. Install with: pip install python-docx"
                    )
            
            # Excel files
            elif file_ext in ['xlsx', 'xls']:
                try:
                    import pandas as pd
                    from io import BytesIO
                    excel_file = pd.ExcelFile(BytesIO(file_content))
                    content_parts = []
                    for sheet_name in excel_file.sheet_names:
                        df = pd.read_excel(BytesIO(file_content), sheet_name=sheet_name)
                        content_parts.append(f"Sheet: {sheet_name}\n{df.to_string(index=False)}")
                    return "\n\n" + "="*50 + "\n\n".join(content_parts)
                except ImportError:
                    logger.warning("pandas not available, cannot process Excel files")
                    raise HTTPException(
                        status_code=400, 
                        detail="Excel processing requires pandas and openpyxl libraries. Install with: pip install pandas openpyxl xlrd"
                    )
            
            # PowerPoint files
            elif file_ext in ['pptx', 'ppt']:
                try:
                    from pptx import Presentation
                    from io import BytesIO
                    prs = Presentation(BytesIO(file_content))
                    content_parts = []
                    for i, slide in enumerate(prs.slides):
                        slide_text = []
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                slide_text.append(shape.text)
                        if slide_text:
                            content_parts.append(f"--- Slide {i + 1} ---\n" + "\n".join(slide_text))
                    return "\n\n".join(content_parts)
                except ImportError:
                    logger.warning("python-pptx not available, cannot process PowerPoint files")
                    raise HTTPException(
                        status_code=400, 
                        detail="PowerPoint processing requires python-pptx library. Install with: pip install python-pptx"
                    )
            
            # RTF files
            elif file_ext == 'rtf':
                try:
                    from striprtf.striprtf import rtf_to_text
                    rtf_content = file_content.decode('utf-8')
                    return rtf_to_text(rtf_content)
                except ImportError:
                    logger.warning("striprtf not available, cannot process RTF files")
                    raise HTTPException(
                        status_code=400, 
                        detail="RTF processing requires striprtf library. Install with: pip install striprtf"
                    )
            
            # LibreOffice formats - basic support
            elif file_ext in ['odt', 'ods', 'odp']:
                logger.warning(f"LibreOffice format {file_ext} requires conversion to supported format")
                raise HTTPException(
                    status_code=400, 
                    detail=f"LibreOffice {file_ext.upper()} files are not directly supported. Please convert to DOCX, PDF, or TXT format."
                )
            
            # Use textract as fallback for other formats if available
            else:
                try:
                    import textract
                    from tempfile import NamedTemporaryFile
                    import os
                    
                    # Create temporary file
                    with NamedTemporaryFile(suffix=f'.{file_ext}', delete=False) as temp_file:
                        temp_file.write(file_content)
                        temp_file.flush()
                        
                        try:
                            # Use textract to extract text
                            extracted_text = textract.process(temp_file.name)
                            return extracted_text.decode('utf-8')
                        finally:
                            # Clean up temporary file
                            os.unlink(temp_file.name)
                            
                except ImportError:
                    logger.warning("textract not available for additional format support")
                    raise HTTPException(
                        status_code=400, 
                        detail=f"Unsupported file type: {filename}. Supported formats: PDF, TXT, MD, CSV, JSON, XML, HTML, DOCX, DOC, XLSX, XLS, PPTX, PPT, RTF. For more formats, install: pip install textract"
                    )
                except Exception as e:
                    logger.error(f"Error processing file with textract: {e}")
                    raise HTTPException(
                        status_code=400, 
                        detail=f"Error processing file {filename}: {str(e)}"
                    )
        
        except HTTPException:
            # Re-raise HTTP exceptions as-is
            raise
        except Exception as e:
            logger.error(f"Unexpected error processing file {filename}: {e}")
            raise HTTPException(
                status_code=500, 
                detail=f"Unexpected error processing file: {str(e)}"
            )
        """Extract text from various file formats supported by LightRAG"""
        try:
            file_ext = filename.lower().split('.')[-1] if '.' in filename else ''
            
            # PDF files
            if file_ext == 'pdf':
                return extract_text_from_pdf_lightrag(file_content)
            
            # Plain text files
            elif file_ext in ['txt', 'md', 'markdown', 'rst']:
                return file_content.decode('utf-8')
            
            # CSV files
            elif file_ext == 'csv':
                return file_content.decode('utf-8')
            
            # JSON files
            elif file_ext == 'json':
                import json
                json_data = json.loads(file_content.decode('utf-8'))
                return json.dumps(json_data, indent=2)
            
            # XML/HTML files
            elif file_ext in ['xml', 'html', 'htm']:
                content = file_content.decode('utf-8')
                try:
                    from bs4 import BeautifulSoup
                    if file_ext in ['html', 'htm']:
                        soup = BeautifulSoup(content, 'html.parser')
                        return soup.get_text(separator='\n', strip=True)
                    else:
                        soup = BeautifulSoup(content, 'xml')
                        return soup.get_text(separator='\n', strip=True)
                except ImportError:
                    # Fallback to raw content if BeautifulSoup not available
                    return content
            
            # Word documents
            elif file_ext in ['docx', 'doc']:
                try:
                    import docx
                    from io import BytesIO
                    doc = docx.Document(BytesIO(file_content))
                    content_parts = []
                    for para in doc.paragraphs:
                        text = para.text.strip()
                        if text:
                            content_parts.append(text)
                    return "\n\n".join(content_parts)
                except ImportError:
                    logger.warning("python-docx not available, cannot process Word documents")
                    raise HTTPException(
                        status_code=400, 
                        detail="Word document processing requires python-docx library"
                    )
            
            # Excel files
            elif file_ext in ['xlsx', 'xls']:
                try:
                    import pandas as pd
                    from io import BytesIO
                    excel_file = pd.ExcelFile(BytesIO(file_content))
                    content_parts = []
                    for sheet_name in excel_file.sheet_names:
                        df = pd.read_excel(BytesIO(file_content), sheet_name=sheet_name)
                        content_parts.append(f"Sheet: {sheet_name}\n{df.to_string(index=False)}")
                    return "\n\n" + "="*50 + "\n\n".join(content_parts)
                except ImportError:
                    logger.warning("pandas not available, cannot process Excel files")
                    raise HTTPException(
                        status_code=400, 
                        detail="Excel processing requires pandas and openpyxl libraries"
                    )
            
            # PowerPoint files
            elif file_ext in ['pptx', 'ppt']:
                try:
                    from pptx import Presentation
                    from io import BytesIO
                    prs = Presentation(BytesIO(file_content))
                    content_parts = []
                    for i, slide in enumerate(prs.slides):
                        slide_text = []
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                slide_text.append(shape.text)
                        if slide_text:
                            content_parts.append(f"--- Slide {i + 1} ---\n" + "\n".join(slide_text))
                    return "\n\n".join(content_parts)
                except ImportError:
                    logger.warning("python-pptx not available, cannot process PowerPoint files")
                    raise HTTPException(
                        status_code=400, 
                        detail="PowerPoint processing requires python-pptx library"
                    )
            
            # RTF files
            elif file_ext == 'rtf':
                try:
                    from striprtf.striprtf import rtf_to_text
                    rtf_content = file_content.decode('utf-8')
                    return rtf_to_text(rtf_content)
                except ImportError:
                    logger.warning("striprtf not available, cannot process RTF files")
                    raise HTTPException(
                        status_code=400, 
                        detail="RTF processing requires striprtf library"
                    )
            
            # LibreOffice formats - basic support
            elif file_ext in ['odt', 'ods', 'odp']:
                logger.warning(f"LibreOffice format {file_ext} requires conversion to supported format")
                raise HTTPException(
                    status_code=400, 
                    detail=f"LibreOffice {file_ext.upper()} files are not directly supported. Please convert to DOCX, PDF, or TXT format."
                )
            
            # Use textract as fallback for other formats if available
            else:
                try:
                    import textract
                    from tempfile import NamedTemporaryFile
                    
                    # Create temporary file
                    with NamedTemporaryFile(suffix=f'.{file_ext}', delete=False) as temp_file:
                        temp_file.write(file_content)
                        temp_file.flush()
                        
                        try:
                            # Use textract to extract text
                            extracted_text = textract.process(temp_file.name)
                            return extracted_text.decode('utf-8')
                        finally:
                            # Clean up temporary file
                            import os
                            os.unlink(temp_file.name)
                            
                except ImportError:
                    logger.warning("textract not available for additional format support")
                    raise HTTPException(
                        status_code=400, 
                        detail=f"Unsupported file type: {filename}. Supported formats: PDF, TXT, MD, CSV, JSON, XML, HTML, DOCX, DOC, XLSX, XLS, PPTX, PPT, RTF"
                    )
                except Exception as e:
                    logger.error(f"Error processing file with textract: {e}")
                    raise HTTPException(
                        status_code=400, 
                        detail=f"Error processing file {filename}: {str(e)}"
                    )
        
        except HTTPException:
            # Re-raise HTTP exceptions as-is
            raise
        except Exception as e:
            logger.error(f"Unexpected error processing file {filename}: {e}")
            raise HTTPException(
                status_code=500, 
                detail=f"Unexpected error processing file: {str(e)}"
            )

    async def create_lightrag_instance(notebook_id: str, llm_provider_config: Dict[str, Any], embedding_provider_config: Dict[str, Any]) -> LightRAG:
        """Create a new LightRAG instance for a notebook with specified provider configurations"""
        working_dir = LIGHTRAG_STORAGE_PATH / notebook_id
        
        # Create directory if it doesn't exist (preserve existing data)
        working_dir.mkdir(exist_ok=True)
        
        try:
            logger.info(f"Creating LightRAG instance for notebook {notebook_id}")
            
            # Get configuration
            llm_provider_type = llm_provider_config.get('type', 'openai')
            llm_model_name = llm_provider_config.get('model', 'gpt-3.5-turbo')
            llm_api_key = llm_provider_config.get('apiKey', '')
            llm_base_url = llm_provider_config.get('baseUrl', '')
            
            embedding_provider_type = embedding_provider_config.get('type', 'openai')
            embedding_model_name = embedding_provider_config.get('model', 'text-embedding-ada-002')
            embedding_api_key = embedding_provider_config.get('apiKey', '')
            embedding_base_url = embedding_provider_config.get('baseUrl', '')
            
            # Log final configuration (auto-detection already done in create_notebook)
            logger.info(f"LLM Provider: {llm_provider_type} - {llm_model_name}")
            logger.info(f"Embedding Provider: {embedding_provider_type} - {embedding_model_name}")
            
            # Remove /v1 suffix from Ollama URLs if present
            if llm_provider_type == 'ollama' and llm_base_url.endswith('/v1'):
                llm_base_url = llm_base_url[:-3]
            if embedding_provider_type == 'ollama' and embedding_base_url.endswith('/v1'):
                embedding_base_url = embedding_base_url[:-3]
            
            # Validate API keys for non-Ollama providers (except Clara Core on port 8091)
            is_clara_core_llm = ':8091' in llm_base_url or llm_base_url.endswith(':8091')
            is_clara_core_embedding = ':8091' in embedding_base_url or embedding_base_url.endswith(':8091')
            
            if llm_provider_type != 'ollama' and not is_clara_core_llm and (not llm_api_key or llm_api_key.strip() == 'your-api-key'):
                raise ValueError(f"Invalid or missing LLM API key for provider type: {llm_provider_type}")
            if embedding_provider_type != 'ollama' and not is_clara_core_embedding and (not embedding_api_key or embedding_api_key.strip() == 'your-api-key'):
                raise ValueError(f"Invalid or missing embedding API key for provider type: {embedding_provider_type}")
            
            # Set Clara Core API key if needed
            if is_clara_core_llm and (not llm_api_key or llm_api_key.strip() == 'your-api-key'):
                llm_api_key = 'claracore'
                logger.info("Using Clara Core LLM - no API key required")
            if is_clara_core_embedding and (not embedding_api_key or embedding_api_key.strip() == 'your-api-key'):
                embedding_api_key = 'claracore'
                logger.info("Using Clara Core embedding - no API key required")
            
            # Determine embedding dimensions based on the embedding model
            embedding_dim = 1536  # Default for OpenAI ada-002
            if 'mxbai-embed-large' in embedding_model_name:
                embedding_dim = 1024
            elif 'e5-large-v2' in embedding_model_name:
                embedding_dim = 1024  # e5-large-v2 produces 1024-dimensional embeddings
            elif 'text-embedding-3-small' in embedding_model_name:
                embedding_dim = 1536
            elif 'text-embedding-3-large' in embedding_model_name:
                embedding_dim = 3072
            elif 'all-MiniLM-L6-v2' in embedding_model_name:
                embedding_dim = 384
            elif 'nomic-embed' in embedding_model_name:
                embedding_dim = 768
            elif 'bge-m3' in embedding_model_name:
                embedding_dim = 1024
            
            logger.info(f"Using embedding dimension: {embedding_dim}")
            
            # Set up LLM function and parameters based on provider type
            if llm_provider_type == 'openai':
                llm_model_func = gpt_4o_mini_complete if 'gpt-4o-mini' in llm_model_name else openai_complete
                llm_model_kwargs = {
                    "api_key": llm_api_key.strip(),
                }
            elif llm_provider_type == 'openai_compatible':
                # Create wrapper function for openai_complete_if_cache
                async def llm_model_func(
                    prompt, system_prompt=None, history_messages=[], keyword_extraction=False, **kwargs
                ) -> str:
                    return await openai_complete_if_cache(
                        llm_model_name,
                        prompt,
                        system_prompt=system_prompt,
                        history_messages=history_messages,
                        api_key=llm_api_key.strip(),
                        base_url=llm_base_url,
                        **kwargs,
                    )
                llm_model_kwargs = {}
            elif llm_provider_type == 'ollama':
                llm_model_func = ollama_model_complete
                llm_model_kwargs = {
                    "host": llm_base_url if llm_base_url else "http://localhost:11434",
                    "options": {"num_ctx": 8192},
                    "timeout": 300,
                }
            else:
                raise ValueError(f"Unsupported LLM provider type: {llm_provider_type}")
            
            # Helper function to batch texts for embedding
            async def batch_embed_texts(texts: list[str], batch_size: int, embed_func):
                """Process texts in batches to avoid size limits"""
                if len(texts) <= batch_size:
                    return await embed_func(texts)
                
                all_embeddings = []
                total_batches = (len(texts) + batch_size - 1) // batch_size
                
                for i in range(0, len(texts), batch_size):
                    batch = texts[i:i + batch_size]
                    batch_num = (i // batch_size) + 1
                    
                    # Log batch details for debugging
                    batch_sizes = [len(text) for text in batch]
                    logger.info(f"Processing embedding batch {batch_num}/{total_batches} ({len(batch)} texts, sizes: {batch_sizes})")
                    
                    try:
                        batch_embeddings = await embed_func(batch)
                        all_embeddings.extend(batch_embeddings)
                        logger.info(f"✅ Batch {batch_num} successful: {len(batch_embeddings)} embeddings")
                        
                        # Small delay between batches to be respectful to the API
                        if i + batch_size < len(texts):
                            await asyncio.sleep(0.5)
                            
                    except Exception as e:
                        error_str = str(e).lower()
                        logger.error(f"❌ Batch {batch_num} failed: {str(e)}")
                        
                        if ('too large' in error_str or 'batch size' in error_str or 
                            'input is too large' in error_str):
                            if batch_size > 1:
                                # Recursively reduce batch size if still too large
                                logger.warning(f"Batch size {batch_size} too large, reducing to {batch_size // 2}")
                                return await batch_embed_texts(texts, batch_size // 2, embed_func)
                            else:
                                # If batch size is 1 and still fails, check if we can chunk the text further
                                if len(batch) == 1 and len(batch[0]) > 200:
                                    logger.warning(f"Individual text too large ({len(batch[0])} chars), attempting to chunk further")
                                    # Split the large text into smaller pieces
                                    large_text = batch[0]
                                    chunks = [large_text[j:j+200] for j in range(0, len(large_text), 180)]  # 20 char overlap
                                    logger.info(f"Split text into {len(chunks)} smaller chunks")
                                    
                                    # Try to embed the chunks individually
                                    chunk_embeddings = []
                                    for chunk in chunks:
                                        try:
                                            chunk_result = await embed_func([chunk])
                                            chunk_embeddings.extend(chunk_result)
                                        except Exception as chunk_error:
                                            logger.error(f"Even small chunk failed ({len(chunk)} chars): {chunk_error}")
                                            raise Exception(f"Text cannot be embedded even after chunking: {str(chunk_error)}")
                                    
                                    all_embeddings.extend(chunk_embeddings)
                                    logger.info(f"✅ Successfully processed large text via chunking: {len(chunk_embeddings)} embeddings")
                                else:
                                    logger.error(f"Individual text too large for embedding: {len(batch[0])} characters")
                                    raise Exception(f"Text too large for embedding (batch size 1 failed): {str(e)}")
                        else:
                            # Re-raise non-batch-size errors
                            raise e
                
                return all_embeddings

            # Set up embedding function based on provider type (following LightRAG documentation pattern)
            if embedding_provider_type == 'openai':
                async def base_openai_embed(texts: list[str]):
                    return await openai_embed(
                        texts,
                        model=embedding_model_name,
                        api_key=embedding_api_key.strip()
                    )
                
                async def embedding_func_lambda(texts: list[str]):
                    # Use smaller batch size for OpenAI to avoid rate limits
                    batch_size = 16 if len(texts) > 16 else len(texts)
                    return await batch_embed_texts(texts, batch_size, base_openai_embed)
                    
            elif embedding_provider_type == 'openai_compatible':
                # Special handling for Clara Core embedding models
                model_to_use = embedding_model_name
                
                async def base_openai_compatible_embed(texts: list[str]):
                    # Retry logic for Clara Core (models need time to load into RAM)
                    max_retries = 5 if is_clara_core_embedding else 2
                    retry_delay = 10 if is_clara_core_embedding else 3  # seconds
                    request_timeout = 180 if is_clara_core_embedding else 60  # seconds
                    
                    for attempt in range(max_retries):
                        try:
                            logger.info(f"Embedding attempt {attempt + 1}/{max_retries} for {model_to_use} ({len(texts)} texts)")
                            
                            # Use asyncio timeout for the request
                            result = await asyncio.wait_for(
                                openai_embed(
                                    texts,
                                    model=model_to_use,
                                    api_key=embedding_api_key.strip(),
                                    base_url=embedding_base_url
                                ),
                                timeout=request_timeout
                            )
                            
                            logger.info(f"Successfully generated {len(result)} embeddings with {model_to_use}")
                            return result
                            
                        except asyncio.TimeoutError:
                            if attempt < max_retries - 1:
                                logger.warning(f"Embedding request timed out (attempt {attempt + 1}/{max_retries}), retrying in {retry_delay}s...")
                                await asyncio.sleep(retry_delay)
                                continue
                            else:
                                error_msg = f"Embedding request timed out after {max_retries} attempts. Clara Core server may need more time to load the model."
                                logger.error(error_msg)
                                raise Exception(error_msg)
                        except Exception as e:
                            error_str = str(e).lower()
                            if attempt < max_retries - 1:
                                # Retry for certain types of errors that might resolve with time
                                if any(keyword in error_str for keyword in ['connection', 'timeout', 'loading', 'model']):
                                    logger.warning(f"Embedding failed (attempt {attempt + 1}/{max_retries}): {e}")
                                    logger.info(f"Retrying in {retry_delay}s... (Clara Core may be loading model)")
                                    await asyncio.sleep(retry_delay)
                                    continue
                                else:
                                    # Check if it's a batch size error
                                    if 'too large' in error_str or 'batch size' in error_str:
                                        logger.warning(f"Batch size error detected: {e}")
                                        raise e  # Let the batching function handle this
                                    # Don't retry for authentication or other permanent errors
                                    raise e
                            else:
                                logger.error(f"Embedding failed after {max_retries} attempts: {e}")
                                raise e

                async def embedding_func_lambda(texts: list[str]):
                    # Use very conservative batch size for Clara Core based on testing
                    # Clara Core with e5-large-v2-q4-0 has strict limits ~500 chars per text
                    if is_clara_core_embedding:
                        # For Clara Core, process each text individually and aggregate chunks
                        final_embeddings = []
                        
                        for text in texts:
                            if len(text) > 400:  # Conservative limit for Clara Core
                                # Split large texts into smaller chunks
                                chunks = [text[i:i+400] for i in range(0, len(text), 350)]  # 50 char overlap
                                logger.info(f"Split large text ({len(text)} chars) into {len(chunks)} chunks for Clara Core")
                                
                                # Get embeddings for all chunks
                                chunk_embeddings = await batch_embed_texts(chunks, 1, base_openai_compatible_embed)
                                
                                # Aggregate chunk embeddings by averaging
                                import numpy as np
                                chunk_array = np.array(chunk_embeddings)
                                aggregated_embedding = np.mean(chunk_array, axis=0).tolist()
                                final_embeddings.append(aggregated_embedding)
                                logger.info(f"Aggregated {len(chunks)} chunk embeddings into single embedding")
                            else:
                                # Small text, process directly
                                single_embedding = await batch_embed_texts([text], 1, base_openai_compatible_embed)
                                final_embeddings.extend(single_embedding)
                        
                        return final_embeddings
                    else:
                        # For other OpenAI-compatible APIs, use larger batch size
                        batch_size = min(16, len(texts)) if len(texts) > 0 else 1
                        return await batch_embed_texts(texts, batch_size, base_openai_compatible_embed)
                    
            elif embedding_provider_type == 'ollama':
                async def base_ollama_embed(texts: list[str]):
                    return await ollama_embed(
                        texts,
                        embed_model=embedding_model_name,
                        host=embedding_base_url if embedding_base_url else "http://localhost:11434"
                    )
                
                async def embedding_func_lambda(texts: list[str]):
                    # Use smaller batch size for Ollama to avoid memory issues
                    batch_size = 4 if len(texts) > 4 else len(texts)
                    return await batch_embed_texts(texts, batch_size, base_ollama_embed)
            else:
                raise ValueError(f"Unsupported embedding provider type: {embedding_provider_type}")
            
            # Determine if using local/ollama providers for optimized configuration
            is_local_llm = llm_provider_type == 'ollama'
            
            # Optimize configuration for local vs remote models
            if is_local_llm:
                # Local/Ollama models: smaller chunks, less aggressive entity extraction
                chunk_token_size = 800
                chunk_overlap_token_size = 80
                entity_extract_max_gleaning = 1
                logger.info(f"Using local model optimization for notebook {notebook_id}")
            else:
                # Remote models: larger chunks, more aggressive entity extraction
                chunk_token_size = 1200
                chunk_overlap_token_size = 100
                entity_extract_max_gleaning = 2
                logger.info(f"Using remote model optimization for notebook {notebook_id}")
            
            # Determine optimal token sizes based on model type
            if 'gemma' in llm_model_name.lower() or 'llama' in llm_model_name.lower():
                # For smaller open source models, use more conservative token limits
                max_token_size = 4096 if 'gemma' in llm_model_name.lower() else 6144
                chunk_token_size = 2000 if is_local_llm else 1200
                chunk_overlap_token_size = 100 if is_local_llm else 100
                entity_extract_max_gleaning = 1  # Reduce complexity for smaller models
                logger.info(f"Using conservative settings for model {llm_model_name}: max_tokens={max_token_size}")
            else:
                # For larger models like GPT-4, use larger token limits
                max_token_size = 8192
                chunk_token_size = chunk_token_size
                chunk_overlap_token_size = chunk_overlap_token_size
                entity_extract_max_gleaning = entity_extract_max_gleaning
                logger.info(f"Using standard settings for model {llm_model_name}: max_tokens={max_token_size}")
            
            # Create LightRAG instance following the official pattern
            logger.info(f"Initializing LightRAG with embedding dimensions: {embedding_dim}")
            rag = LightRAG(
                working_dir=str(working_dir),
                llm_model_func=llm_model_func,
                llm_model_name=llm_model_name,
                llm_model_max_token_size=max_token_size,
                llm_model_kwargs=llm_model_kwargs,
                embedding_func=EmbeddingFunc(
                    embedding_dim=embedding_dim,
                    max_token_size=max_token_size,
                    func=embedding_func_lambda,
                ),
                chunk_token_size=chunk_token_size,
                chunk_overlap_token_size=chunk_overlap_token_size,
                entity_extract_max_gleaning=entity_extract_max_gleaning,
            )
            
            # Initialize storages
            try:
                await rag.initialize_storages()
                await initialize_pipeline_status()
                logger.info(f"LightRAG storages initialized for notebook {notebook_id}")
            except Exception as init_error:
                logger.warning(f"Storage initialization error (may be expected): {init_error}")
            
            logger.info(f"Successfully created LightRAG instance for notebook {notebook_id}")
            return rag
            
        except ValueError as ve:
            # Handle configuration errors
            logger.error(f"Configuration error for notebook {notebook_id}: {ve}")
            raise HTTPException(status_code=400, detail=f"Configuration error: {str(ve)}")
        except Exception as e:
            logger.error(f"Error creating LightRAG instance for notebook {notebook_id}: {e}")
            logger.error(f"Full traceback: {traceback.format_exc()}")
            raise HTTPException(status_code=500, detail=f"Error initializing RAG system: {str(e)}")

    async def get_lightrag_instance(notebook_id: str) -> LightRAG:
        """Get or create LightRAG instance for a notebook"""
        if notebook_id not in lightrag_instances:
            if notebook_id not in lightrag_notebooks_db:
                raise HTTPException(status_code=404, detail="Notebook not found")
            
            lightrag_instances[notebook_id] = await create_lightrag_instance(notebook_id, lightrag_notebooks_db[notebook_id]["llm_provider"], lightrag_notebooks_db[notebook_id]["embedding_provider"])
        
        return lightrag_instances[notebook_id]

    def auto_detect_provider_type(provider_config: Dict[str, Any]) -> Dict[str, Any]:
        """Auto-detect provider type based on baseUrl and return updated config"""
        provider_config = provider_config.copy()  # Don't modify original
        provider_type = provider_config.get('type', 'openai')
        base_url = provider_config.get('baseUrl', '')
        
        # Auto-detect for non-ollama providers
        if provider_type != 'ollama':
            if base_url in ['https://api.openai.com/v1', 'https://api.openai.com']:
                provider_config['type'] = 'openai'
                logger.info(f"Auto-detected provider type as 'openai' based on baseUrl: {base_url}")
            else:
                provider_config['type'] = 'openai_compatible'
                logger.info(f"Auto-detected provider type as 'openai_compatible' based on baseUrl: {base_url}")
        
        return provider_config

    def validate_notebook_exists(notebook_id: str):
        """Validate that a notebook exists"""
        if notebook_id not in lightrag_notebooks_db:
            raise HTTPException(status_code=404, detail="Notebook not found")

    async def process_notebook_document_with_delay(notebook_id: str, document_id: str, text_content: str, delay_seconds: int):
        """Wrapper to add delay before processing document"""
        if delay_seconds > 0:
            await asyncio.sleep(delay_seconds)
        
        await process_notebook_document(notebook_id, document_id, text_content)

    async def process_notebook_document(notebook_id: str, document_id: str, text_content: str):
        """Background task to process document with LightRAG"""
        try:
            # Validate inputs
            if not text_content or not text_content.strip():
                raise ValueError("Document content is empty")
            
            logger.info(f"Starting document processing for {document_id} in notebook {notebook_id}")
            
            rag = await get_lightrag_instance(notebook_id)
            
            # Get notebook data to check provider type
            notebook_data = lightrag_notebooks_db[notebook_id]
            llm_provider = notebook_data.get("llm_provider", {})
            llm_provider_type = llm_provider.get("type", "openai")
            is_local_provider = llm_provider_type == 'ollama'
            
            # Adjust content size and timeout based on provider type
            if is_local_provider:
                # For local/ollama models, use smaller chunks and longer timeout
                max_content_size = 300000  # 300KB of text for local models
                processing_timeout = 1800.0  # 30 minutes for local/ollama
                logger.info(f"Using local provider settings: max_size={max_content_size}, timeout={processing_timeout}s")
            else:
                # For remote models, use larger chunks and shorter timeout  
                max_content_size = 800000  # 800KB of text for remote models
                processing_timeout = 900.0  # 15 minutes for remote models
                logger.info(f"Using remote provider settings: max_size={max_content_size}, timeout={processing_timeout}s")
            
            if len(text_content) > max_content_size:
                logger.warning(f"Document {document_id} is very large ({len(text_content)} chars), truncating to {max_content_size}")
                text_content = text_content[:max_content_size] + "\n\n[Content truncated due to size limits]"
            
            # Create a more specific document ID to avoid conflicts
            import hashlib
            import time
            timestamp = str(int(time.time() * 1000))  # milliseconds
            content_hash = hashlib.md5(text_content.encode()).hexdigest()[:8]
            prefixed_doc_id = f"doc_{notebook_id}_{document_id}_{timestamp}_{content_hash}"
            
            logger.info(f"Processing document {document_id} ({len(text_content)} chars) with ID {prefixed_doc_id}")
            
            # Set document status to processing before starting
            if document_id in lightrag_documents_db:
                lightrag_documents_db[document_id]["status"] = "processing"
                lightrag_documents_db[document_id]["processed_at"] = datetime.now()
                save_documents_db()
            
            # Get document metadata including file path for citations
            document_data = lightrag_documents_db[document_id]
            filename = document_data["filename"]
            file_path = document_data.get("file_path", f"documents/{filename}")
            
            # Insert document text into LightRAG
            try:
                logger.info(f"Starting LightRAG insertion for document {document_id}")
                
                # Use asyncio timeout to prevent hanging
                await asyncio.wait_for(
                        rag.ainsert(text_content, ids=[prefixed_doc_id]),
                        timeout=processing_timeout
                )
                
                logger.info(f"Successfully inserted document {document_id} into LightRAG")
                
            except asyncio.TimeoutError:
                timeout_minutes = int(processing_timeout / 60)
                error_msg = f"Document processing timed out after {timeout_minutes} minutes. This can happen with complex documents and local models. Try using a more powerful model or splitting the document into smaller parts."
                logger.error(error_msg)
                raise Exception(error_msg)
            except Exception as insert_error:
                logger.error(f"LightRAG insertion failed for document {document_id}: {insert_error}")
                
                # Check for common errors and provide helpful messages
                error_str = str(insert_error).lower()
                if "connection" in error_str or "timeout" in error_str:
                    raise Exception(f"Connection error during document processing: {str(insert_error)}. Please check your provider configuration and network connection.")
                elif "api key" in error_str or "unauthorized" in error_str:
                    raise Exception(f"Authentication error: {str(insert_error)}. Please check your API key configuration.")
                elif "model" in error_str and "not found" in error_str:
                    raise Exception(f"Model not found: {str(insert_error)}. Please check your model name configuration.")
                else:
                    raise Exception(f"Document processing failed: {str(insert_error)}")
            
            # Force cache clear after inserting document
            try:
                await asyncio.wait_for(rag.aclear_cache(), timeout=60.0)
                logger.info(f"Cache cleared for document {document_id}")
            except asyncio.TimeoutError:
                logger.warning("Cache clear timed out, continuing anyway")
            except Exception as cache_error:
                logger.warning(f"Cache clear failed: {cache_error}, continuing anyway")
            
            # Update document status to completed
            if document_id in lightrag_documents_db:
                lightrag_documents_db[document_id]["status"] = "completed"
                lightrag_documents_db[document_id]["lightrag_id"] = prefixed_doc_id
                lightrag_documents_db[document_id]["completed_at"] = datetime.now()
                # Clear any previous error
                if "error" in lightrag_documents_db[document_id]:
                    del lightrag_documents_db[document_id]["error"]
                
                # Optional: Clear content after successful processing to save space
                # Keep content for failed documents so they can be retried
                # For completed documents, the content is already in LightRAG
                if "content" in lightrag_documents_db[document_id]:
                    content_size = len(lightrag_documents_db[document_id]["content"])
                    del lightrag_documents_db[document_id]["content"]
                    logger.info(f"Cleared content ({content_size} chars) for completed document {document_id}")
                
                # Also clean up content file if it exists (document is now safely in LightRAG)
                if "content_file" in lightrag_documents_db[document_id]:
                    try:
                        content_file = Path(lightrag_documents_db[document_id]["content_file"])
                        if content_file.exists():
                            content_file.unlink()
                            logger.info(f"Cleaned up content file: {content_file}")
                        del lightrag_documents_db[document_id]["content_file"]
                    except Exception as e:
                        logger.warning(f"Failed to clean up content file: {e}")
            
            # Clear summary cache since a new document has been processed
            if notebook_id in lightrag_notebooks_db:
                if "summary_cache" in lightrag_notebooks_db[notebook_id]:
                    del lightrag_notebooks_db[notebook_id]["summary_cache"]
                    logger.info(f"Cleared summary cache for notebook {notebook_id}")
                if "docs_fingerprint" in lightrag_notebooks_db[notebook_id]:
                    del lightrag_notebooks_db[notebook_id]["docs_fingerprint"]
            
            # Save changes to disk
            save_documents_db()
            save_notebooks_db()
            
            logger.info(f"Successfully completed processing document {document_id} in notebook {notebook_id}")
            
        except Exception as e:
            error_msg = str(e)
            logger.error(f"Error processing document {document_id} in notebook {notebook_id}: {error_msg}")
            logger.error(f"Full error traceback: {traceback.format_exc()}")
            
            # Update document status to failed
            if document_id in lightrag_documents_db:
                lightrag_documents_db[document_id]["status"] = "failed"
                lightrag_documents_db[document_id]["error"] = error_msg
                lightrag_documents_db[document_id]["failed_at"] = datetime.now()
                # Save changes to disk even on failure
                save_documents_db()

class TTSRequest(BaseModel):
    text: str
    language: Optional[str] = "en"
    engine: Optional[str] = None  # "gtts", "pyttsx3", "kokoro", "kokoro-onnx", or None for current engine
    slow: Optional[bool] = False
    voice: Optional[str] = "af_sarah"  # Voice for Kokoro engines
    speed: Optional[float] = 1.0  # Speed for Kokoro engines (0.5-2.0)

# LightRAG Data Models
class NotebookCreate(BaseModel):
    name: str = Field(..., description="Name of the notebook")
    description: Optional[str] = Field(None, description="Description of the notebook")
    # Add provider configuration
    llm_provider: Dict[str, Any] = Field(..., description="LLM provider configuration")
    embedding_provider: Dict[str, Any] = Field(..., description="Embedding provider configuration")

class NotebookResponse(BaseModel):
    id: str = Field(..., description="Notebook ID")
    name: str = Field(..., description="Notebook name")
    description: Optional[str] = Field(None, description="Notebook description")
    created_at: datetime = Field(..., description="Creation timestamp")
    document_count: int = Field(0, description="Number of documents in notebook")
    # Add provider information to response (optional for backward compatibility)
    llm_provider: Optional[Dict[str, Any]] = Field(None, description="LLM provider configuration")
    embedding_provider: Optional[Dict[str, Any]] = Field(None, description="Embedding provider configuration")

class NotebookDocumentResponse(BaseModel):
    id: str = Field(..., description="Document ID")
    filename: str = Field(..., description="Original filename")
    notebook_id: str = Field(..., description="Notebook ID")
    uploaded_at: datetime = Field(..., description="Upload timestamp")
    status: str = Field(..., description="Processing status")
    error: Optional[str] = Field(None, description="Error message if processing failed")
    file_path: Optional[str] = Field(None, description="File path for citation tracking")

class NotebookQueryRequest(BaseModel):
    question: str = Field(..., description="Question to ask")
    mode: str = Field("hybrid", description="Query mode: local, global, hybrid, naive, mix")
    response_type: str = Field("Multiple Paragraphs", description="Response format")
    top_k: int = Field(60, description="Number of top items to retrieve")
    # Add optional provider override for query
    llm_provider: Optional[Dict[str, Any]] = Field(None, description="Override LLM provider for this query")
    # Add chat history support
    use_chat_history: bool = Field(True, description="Whether to use chat history for context")

class NotebookQueryResponse(BaseModel):
    answer: str = Field(..., description="Generated answer")
    mode: str = Field(..., description="Query mode used")
    context_used: bool = Field(True, description="Whether context was used")
    citations: Optional[List[Dict[str, Any]]] = Field(None, description="Citation information with sources")
    # Enhanced citation support
    source_documents: Optional[List[Dict[str, Any]]] = Field(None, description="Source documents with highlighted passages")
    chat_context_used: bool = Field(False, description="Whether previous chat history was used")

class ChatMessage(BaseModel):
    role: str = Field(..., description="Message role: user or assistant")
    content: str = Field(..., description="Message content")
    timestamp: datetime = Field(default_factory=datetime.now, description="Message timestamp")
    citations: Optional[List[Dict[str, Any]]] = Field(None, description="Citations for assistant messages")

class ChatHistoryResponse(BaseModel):
    notebook_id: str = Field(..., description="Notebook ID")
    messages: List[ChatMessage] = Field(..., description="Chat messages")
    total_messages: int = Field(..., description="Total number of messages")

class DocumentSummaryRequest(BaseModel):
    include_details: bool = Field(True, description="Include document-level details")
    max_length: str = Field("medium", description="Summary length: short, medium, long")

class QueryTemplate(BaseModel):
    id: str = Field(..., description="Template ID")
    name: str = Field(..., description="Template name")
    description: str = Field(..., description="Template description")
    question_template: str = Field(..., description="Question template with placeholders")
    category: str = Field(..., description="Template category")
    use_case: str = Field(..., description="When to use this template")

class DocumentRetryResponse(BaseModel):
    message: str = Field(..., description="Success message")
    document_id: str = Field(..., description="Document ID that was retried")
    status: str = Field(..., description="New document status after retry initiation")

@app.get("/")
def read_root():
    """Root endpoint for basic health check"""
    return {
        "status": "ok", 
        "service": "Clara Backend", 
        "port": PORT,
        "uptime": str(datetime.now() - datetime.fromisoformat(START_TIME)),
        "start_time": START_TIME
    }

@app.get("/health")
def health_check():
    """Health check endpoint"""
    return {
        "status": "healthy",
        "port": PORT,
        "uptime": str(datetime.now() - datetime.fromisoformat(START_TIME))
    }

# LightRAG Notebook endpoints
if LIGHTRAG_AVAILABLE:
    @app.post("/notebooks", response_model=NotebookResponse)
    async def create_notebook(notebook: NotebookCreate):
        """Create a new notebook using LightRAG"""
        notebook_id = str(uuid.uuid4())
        
        # Add debugging logging
        logger.info(f"Creating notebook with data: name={notebook.name}, description={notebook.description}")
        logger.info(f"Original LLM Provider: {notebook.llm_provider}")
        logger.info(f"Original Embedding Provider: {notebook.embedding_provider}")
        
        # Auto-detect provider types before saving
        corrected_llm_provider = auto_detect_provider_type(notebook.llm_provider)
        corrected_embedding_provider = auto_detect_provider_type(notebook.embedding_provider)
        
        logger.info(f"Corrected LLM Provider: {corrected_llm_provider}")
        logger.info(f"Corrected Embedding Provider: {corrected_embedding_provider}")
        
        notebook_data = {
            "id": notebook_id,
            "name": notebook.name,
            "description": notebook.description,
            "created_at": datetime.now(),
            "document_count": 0,
            "llm_provider": corrected_llm_provider,
            "embedding_provider": corrected_embedding_provider
        }
        
        # Log the notebook data before saving
        logger.info(f"Notebook data before saving: {notebook_data}")
        
        lightrag_notebooks_db[notebook_id] = notebook_data
        
        # Create LightRAG instance for this notebook
        try:
            await create_lightrag_instance(notebook_id, corrected_llm_provider, corrected_embedding_provider)
            logger.info(f"Created notebook {notebook_id}: {notebook.name}")
            # Save to disk after successful creation
            save_notebooks_db()
            
            # Log the saved data
            logger.info(f"Saved notebook data: {lightrag_notebooks_db[notebook_id]}")
        except Exception as e:
            # Clean up if LightRAG creation fails
            del lightrag_notebooks_db[notebook_id]
            raise
        
        # Create response and log it
        response = NotebookResponse(**notebook_data)
        logger.info(f"Response being returned: {response.model_dump()}")
        
        return response

    @app.get("/notebooks", response_model=List[NotebookResponse])
    async def list_notebooks():
        """List all notebooks"""
        logger.info(f"Listing notebooks, found {len(lightrag_notebooks_db)} notebooks")
        notebooks = []
        for notebook_id, notebook in lightrag_notebooks_db.items():
            logger.info(f"Processing notebook {notebook_id}: {notebook}")
            
            # Create a copy to avoid modifying the original
            notebook_copy = notebook.copy()
            
            # Add default provider configuration if missing (for backward compatibility)
            if "llm_provider" not in notebook_copy:
                logger.info(f"Adding default LLM provider for notebook {notebook_id}")
                notebook_copy["llm_provider"] = {
                    "name": "OpenAI",
                    "type": "openai",
                    "baseUrl": "https://api.openai.com/v1",
                    "apiKey": "your-api-key",
                    "model": "gpt-4o-mini"
                }
            if "embedding_provider" not in notebook_copy:
                logger.info(f"Adding default embedding provider for notebook {notebook_id}")
                notebook_copy["embedding_provider"] = {
                    "name": "OpenAI",
                    "type": "openai", 
                    "baseUrl": "https://api.openai.com/v1",
                    "apiKey": "your-api-key",
                    "model": "text-embedding-ada-002"
                }
            
            notebook_response = NotebookResponse(**notebook_copy)
            logger.info(f"Notebook response for {notebook_id}: {notebook_response.model_dump()}")
            notebooks.append(notebook_response)
        return notebooks

    @app.get("/notebooks/{notebook_id}", response_model=NotebookResponse)
    async def get_notebook(notebook_id: str):
        """Get a specific notebook"""
        validate_notebook_exists(notebook_id)
        notebook = lightrag_notebooks_db[notebook_id].copy()
        
        # Add default provider configuration if missing (for backward compatibility)
        if "llm_provider" not in notebook:
            notebook["llm_provider"] = {
                "name": "OpenAI",
                "type": "openai",
                "baseUrl": "https://api.openai.com/v1",
                "apiKey": "your-api-key",
                "model": "gpt-4o-mini"
            }
        if "embedding_provider" not in notebook:
            notebook["embedding_provider"] = {
                "name": "OpenAI",
                "type": "openai", 
                "baseUrl": "https://api.openai.com/v1",
                "apiKey": "your-api-key",
                "model": "text-embedding-ada-002"
            }
        
        return NotebookResponse(**notebook)

    @app.delete("/notebooks/{notebook_id}")
    async def delete_notebook(notebook_id: str):
        """Delete a notebook and all its documents"""
        validate_notebook_exists(notebook_id)
        
        # Remove all documents from this notebook
        notebook_docs = [doc_id for doc_id, doc in lightrag_documents_db.items() 
                        if doc["notebook_id"] == notebook_id]
        
        for doc_id in notebook_docs:
            del lightrag_documents_db[doc_id]
        
        # Remove LightRAG instance
        if notebook_id in lightrag_instances:
            del lightrag_instances[notebook_id]
        
        # Remove notebook
        del lightrag_notebooks_db[notebook_id]
        
        # Clean up storage directory
        storage_dir = LIGHTRAG_STORAGE_PATH / notebook_id
        if storage_dir.exists():
            shutil.rmtree(storage_dir, ignore_errors=True)
        
        # Save changes to disk
        save_notebooks_db()
        save_documents_db()
        
        logger.info(f"Deleted notebook {notebook_id}")
        return {"message": "Notebook deleted successfully"}

    @app.post("/notebooks/{notebook_id}/documents", response_model=List[NotebookDocumentResponse])
    async def upload_notebook_documents(
        notebook_id: str,
        background_tasks: BackgroundTasks,
        files: List[UploadFile] = File(...)
    ):
        """Upload multiple documents to a notebook"""
        validate_notebook_exists(notebook_id)
        
        if not files:
            raise HTTPException(status_code=400, detail="No files provided")
        
        uploaded_documents = []
        
        # Process files sequentially to avoid conflicts
        for i, file in enumerate(files):
            if not file.filename:
                continue
                
            # Generate document ID
            document_id = str(uuid.uuid4())
            
            # Read file content
            try:
                file_content = await file.read()
                
                # Extract text based on file type
                text_content = await extract_text_from_file(file.filename, file_content)
                
                # Validate text content
                if not text_content.strip():
                    logger.warning(f"File {file.filename} appears to be empty")
                    continue
                
                # Create file path for citation tracking
                file_path = f"notebooks/{notebook_id}/{file.filename}"
                
                # Create document record with content stored for potential retry
                document_data = {
                    "id": document_id,
                    "filename": file.filename,
                    "notebook_id": notebook_id,
                    "uploaded_at": datetime.now(),
                    "status": "processing",
                    "file_path": file_path,
                    "content": text_content  # Store content for retry functionality
                }
                
                # Debug: Log that content is being stored
                content_length = len(text_content)
                logger.info(f"Storing document {document_id} with content length: {content_length} characters")
                
                # Add content size info for monitoring
                document_data["content_size"] = content_length
                
                # For very large documents, consider storing content in a separate file
                # This prevents JSON serialization issues with very large documents
                if content_length > 1000000:  # 1MB threshold
                    logger.info(f"Large document detected ({content_length} chars), using file storage for content")
                    content_file = Path(data_dir) / f"content_{document_id}.txt"
                    try:
                        with open(content_file, 'w', encoding='utf-8') as f:
                            f.write(text_content)
                        document_data["content_file"] = str(content_file)
                        # Still store content in memory for immediate use, but fallback exists
                        logger.info(f"Content backed up to {content_file}")
                    except Exception as e:
                        logger.warning(f"Failed to create content backup file: {e}")
                
                lightrag_documents_db[document_id] = document_data
                
                # Add background task for document processing with a slight delay to avoid conflicts
                # Increase delay for larger documents or more concurrent uploads
                delay_seconds = min(i * 3, 30)  # 3 second delay between docs, max 30 seconds
                background_tasks.add_task(
                    process_notebook_document_with_delay, 
                    notebook_id, 
                    document_id, 
                    text_content,
                    delay_seconds
                )
                
                # Update notebook document count
                lightrag_notebooks_db[notebook_id]["document_count"] += 1
                
                uploaded_documents.append(NotebookDocumentResponse(**document_data))
                
            except Exception as e:
                logger.error(f"Error processing file {file.filename}: {e}")
                raise HTTPException(
                    status_code=500, 
                    detail=f"Error processing file {file.filename}: {str(e)}"
                )
        
        # Save changes to disk after all uploads
        save_documents_db()
        save_notebooks_db()
        
        logger.info(f"Uploaded {len(uploaded_documents)} documents to notebook {notebook_id}")
        return uploaded_documents

    @app.get("/notebooks/{notebook_id}/documents", response_model=List[NotebookDocumentResponse])
    async def list_notebook_documents(notebook_id: str):
        """List all documents in a notebook"""
        validate_notebook_exists(notebook_id)
        
        notebook_documents = [
            NotebookDocumentResponse(**doc) 
            for doc in lightrag_documents_db.values() 
            if doc["notebook_id"] == notebook_id
        ]
        
        return notebook_documents

    @app.delete("/notebooks/{notebook_id}/documents/{document_id}")
    async def delete_notebook_document(notebook_id: str, document_id: str):
        """Delete a specific document from a notebook"""
        validate_notebook_exists(notebook_id)
        
        if document_id not in lightrag_documents_db:
            raise HTTPException(status_code=404, detail="Document not found")
        
        if lightrag_documents_db[document_id]["notebook_id"] != notebook_id:
            raise HTTPException(status_code=400, detail="Document does not belong to this notebook")
        
        try:
            # Remove from LightRAG using the stored LightRAG ID
            rag = await get_lightrag_instance(notebook_id)
            document_data = lightrag_documents_db[document_id]
            
            # Use the stored LightRAG ID if available, otherwise construct it
            lightrag_id = document_data.get("lightrag_id", f"doc_{notebook_id}_{document_id}")
            await rag.adelete_by_doc_id(lightrag_id)
            
            # Clear cache after deleting document
            await rag.aclear_cache()
            
            # Clean up content file if it exists
            if "content_file" in document_data:
                try:
                    content_file = Path(document_data["content_file"])
                    if content_file.exists():
                        content_file.unlink()
                        logger.info(f"Cleaned up content file during deletion: {content_file}")
                except Exception as e:
                    logger.warning(f"Failed to clean up content file during deletion: {e}")
            
            # Remove from database
            del lightrag_documents_db[document_id]
            
            # Update notebook document count
            lightrag_notebooks_db[notebook_id]["document_count"] -= 1
            
            # Clear summary cache since documents have changed
            if "summary_cache" in lightrag_notebooks_db[notebook_id]:
                del lightrag_notebooks_db[notebook_id]["summary_cache"]
            if "docs_fingerprint" in lightrag_notebooks_db[notebook_id]:
                del lightrag_notebooks_db[notebook_id]["docs_fingerprint"]
            
            # Save changes to disk
            save_documents_db()
            save_notebooks_db()
            
            logger.info(f"Deleted document {document_id} from notebook {notebook_id}")
            return {"message": "Document deleted successfully"}
            
        except Exception as e:
            logger.error(f"Error deleting document {document_id}: {e}")
            raise HTTPException(status_code=500, detail=f"Error deleting document: {str(e)}")

    @app.post("/notebooks/{notebook_id}/documents/{document_id}/retry", response_model=DocumentRetryResponse)
    async def retry_failed_document(notebook_id: str, document_id: str, background_tasks: BackgroundTasks):
        """Retry processing a failed document"""
        validate_notebook_exists(notebook_id)
        
        if document_id not in lightrag_documents_db:
            raise HTTPException(status_code=404, detail="Document not found")
        
        document_data = lightrag_documents_db[document_id]
        
        if document_data["notebook_id"] != notebook_id:
            raise HTTPException(status_code=400, detail="Document does not belong to this notebook")
        
        # Check if document is in a retryable state
        if document_data["status"] not in ["failed"]:
            raise HTTPException(
                status_code=400, 
                detail=f"Document is in '{document_data['status']}' state. Only 'failed' documents can be retried."
            )
        
        try:
            logger.info(f"Retrying document {document_id} in notebook {notebook_id}")
            
            # Debug: Log document data keys
            logger.info(f"Document data keys: {list(document_data.keys())}")
            
            # Get the original text content from the failed document
            # Check if we have stored content or need to re-extract it
            text_content = document_data.get("content")
            
            # If content is not in memory, try to load from content file (for large documents)
            if not text_content and "content_file" in document_data:
                try:
                    content_file = Path(document_data["content_file"])
                    if content_file.exists():
                        with open(content_file, 'r', encoding='utf-8') as f:
                            text_content = f.read()
                        logger.info(f"Loaded content from backup file: {content_file}")
                    else:
                        logger.warning(f"Content file not found: {content_file}")
                except Exception as e:
                    logger.error(f"Failed to load content from file: {e}")
            
            # Debug: Log content availability
            if text_content:
                logger.info(f"Content available for retry, length: {len(text_content)} characters")
            else:
                logger.warning(f"No content found in document data. Available keys: {list(document_data.keys())}")
            
            if not text_content:
                # If content wasn't stored, we need the original file to retry
                raise HTTPException(
                    status_code=400, 
                    detail="Original document content not available for retry. Please re-upload the document."
                )
            
            # Reset document status to processing
            document_data["status"] = "processing"
            document_data["processed_at"] = datetime.now()
            
            # Clear previous error information
            if "failed_at" in document_data:
                del document_data["failed_at"]
            if "error_message" in document_data:
                del document_data["error_message"]
            if "error_details" in document_data:
                del document_data["error_details"]
            
            # Update the document in database
            lightrag_documents_db[document_id] = document_data
            save_documents_db()
            
            # Add background task to retry processing
            # The LightRAG cache will automatically skip chunks that were already processed
            background_tasks.add_task(
                process_notebook_document, 
                notebook_id, 
                document_id, 
                text_content
            )
            
            logger.info(f"Retry initiated for document {document_id}")
            return DocumentRetryResponse(
                message="Document retry initiated successfully",
                document_id=document_id,
                status="processing"
            )
            
        except HTTPException:
            # Re-raise HTTP exceptions as-is
            raise
        except Exception as e:
            logger.error(f"Error retrying document {document_id}: {e}")
            # Reset status back to failed if retry setup failed
            document_data["status"] = "failed"
            lightrag_documents_db[document_id] = document_data
            save_documents_db()
            raise HTTPException(status_code=500, detail=f"Error initiating retry: {str(e)}")

    @app.post("/notebooks/{notebook_id}/query", response_model=NotebookQueryResponse)
    async def query_notebook(notebook_id: str, query: NotebookQueryRequest):
        """Query a notebook with a question"""
        validate_notebook_exists(notebook_id)
        
        try:
            logger.info(f"Query request for notebook {notebook_id}")
            
            # Get the current RAG instance or create a new one if provider is overridden
            if query.llm_provider:
                # Use override provider for this query
                logger.info(f"Using override LLM provider for query: {query.llm_provider.get('name', 'Unknown')}")
                notebook = lightrag_notebooks_db[notebook_id]
                embedding_provider = notebook["embedding_provider"]  # Keep existing embedding provider
                
                # Create temporary RAG instance with override provider
                rag = await create_lightrag_instance(
                    f"{notebook_id}_temp",
                    query.llm_provider,
                    embedding_provider
                )
            else:
                # Use existing RAG instance
                rag = await get_lightrag_instance(notebook_id)
            
            # Get notebook and model information for query optimization
            notebook = lightrag_notebooks_db[notebook_id]
            llm_provider = query.llm_provider or notebook.get("llm_provider", {})
            model_name = llm_provider.get("model", "").lower()
            
            # Adjust query parameters based on model capabilities
            adjusted_top_k = query.top_k
            adjusted_mode = query.mode
            
            # For smaller models like Gemma, use more conservative query parameters
            if 'gemma' in model_name or ('llama' in model_name and any(size in model_name for size in ['3b', '4b', '7b'])):
                logger.info(f"Optimizing query for smaller model: {model_name}")
                # Reduce top_k to limit context size
                adjusted_top_k = min(query.top_k, 30)
                # For very complex queries, prefer local mode to reduce context
                if query.mode == "global" and len(query.question.split()) > 20:
                    adjusted_mode = "hybrid"
                    logger.info("Switching from global to hybrid mode for complex query on small model")
            
            # Create query parameters
            query_param = QueryParam(
                mode=adjusted_mode,
                response_type=query.response_type,
                top_k=adjusted_top_k,
            )
            
            # Perform query with fallback handling for context size issues
            try:
                result = await rag.aquery(query.question, param=query_param)
            except Exception as query_error:
                error_str = str(query_error).lower()
                
                # Check if it's a context size error
                if any(keyword in error_str for keyword in ['context size', 'context length', 'token limit', 'exceeds', 'too long']):
                    logger.warning(f"Context size error detected, attempting recovery: {query_error}")
                    
                    # Try with more aggressive reduction
                    if adjusted_mode == "global":
                        logger.info("Retrying with local mode instead of global")
                        fallback_param = QueryParam(
                            mode="local",
                            response_type=query.response_type,
                            top_k=min(20, adjusted_top_k),
                        )
                        result = await rag.aquery(query.question, param=fallback_param)
                        adjusted_mode = "local"
                    elif adjusted_mode == "hybrid":
                        logger.info("Retrying with naive mode instead of hybrid")
                        fallback_param = QueryParam(
                            mode="naive",
                            response_type=query.response_type,
                            top_k=min(15, adjusted_top_k),
                        )
                        result = await rag.aquery(query.question, param=fallback_param)
                        adjusted_mode = "naive"
                    else:
                        # Already using simplest mode, try with minimal context
                        logger.info("Retrying with minimal context")
                        fallback_param = QueryParam(
                            mode="naive",
                            response_type="Single Paragraph",
                            top_k=5,
                        )
                        result = await rag.aquery(query.question, param=fallback_param)
                        adjusted_mode = "naive"
                else:
                    # Re-raise non-context-size errors
                    raise query_error
            
            # Extract citation information from the result if available
            citations = []
            try:
                # Check if the result contains citation information
                # LightRAG may return metadata about sources used
                # For now, we'll extract from document metadata
                notebook_documents = [
                    doc for doc in lightrag_documents_db.values() 
                    if doc["notebook_id"] == notebook_id and doc["status"] == "completed"
                ]
                
                # Create citations list with available document information
                for doc in notebook_documents:
                    citation = {
                        "filename": doc["filename"],
                        "file_path": doc.get("file_path", f"documents/{doc['filename']}"),
                        "document_id": doc["id"],
                        "title": doc["filename"].replace('_', ' ').replace('.txt', '').replace('.pdf', '').replace('.md', '').title()
                    }
                    citations.append(citation)
                
                # Limit citations to prevent overwhelming the response
                citations = citations[:10] if citations else None
                
            except Exception as citation_error:
                logger.warning(f"Error extracting citations: {citation_error}")
                citations = None
            
            return NotebookQueryResponse(
                answer=result,
                mode=adjusted_mode,
                context_used=True,
                citations=citations,
                source_documents=None,  # Will be enhanced later
                chat_context_used=False  # Will be enhanced when chat history is implemented
            )
            
        except Exception as e:
            logger.error(f"Error processing query for notebook {notebook_id}: {e}")
            raise HTTPException(status_code=500, detail=f"Error processing query: {str(e)}")

    @app.post("/notebooks/{notebook_id}/summary", response_model=NotebookQueryResponse)
    async def generate_notebook_summary(notebook_id: str):
        """Generate an automatic summary of all documents in the notebook"""
        validate_notebook_exists(notebook_id)
        
        try:
            logger.info(f"Summary generation request for notebook {notebook_id}")
            
            # Check if there are any completed documents
            notebook_documents = [
                doc for doc in lightrag_documents_db.values() 
                if doc["notebook_id"] == notebook_id and doc["status"] == "completed"
            ]
            
            if not notebook_documents:
                return NotebookQueryResponse(
                    answer="No documents have been processed yet. Please upload and wait for documents to be processed before generating a summary.",
                    mode="hybrid",
                    context_used=False,
                    citations=None,
                    source_documents=None,
                    chat_context_used=False
                )
            
            # Create a fingerprint of current documents (using document IDs and upload times)
            current_docs_fingerprint = "|".join(sorted([
                f"{doc['id']}:{doc['uploaded_at'].isoformat() if isinstance(doc['uploaded_at'], datetime) else doc['uploaded_at']}" 
                for doc in notebook_documents
            ]))
            
            # Check if we have a cached summary that's still valid
            notebook_data = lightrag_notebooks_db[notebook_id]
            cached_summary = notebook_data.get("summary_cache")
            cached_fingerprint = notebook_data.get("docs_fingerprint")
            
            # If we have a valid cached summary, return it
            if (cached_summary and cached_fingerprint and 
                cached_fingerprint == current_docs_fingerprint):
                logger.info(f"Returning cached summary for notebook {notebook_id}")
                
                # Extract citation information for all completed documents
                citations = []
                try:
                    for doc in notebook_documents:
                        citation = {
                            "filename": doc["filename"],
                            "file_path": doc.get("file_path", f"documents/{doc['filename']}"),
                            "document_id": doc["id"],
                            "title": doc["filename"].replace('_', ' ').replace('.txt', '').replace('.pdf', '').replace('.md', '').title()
                        }
                        citations.append(citation)
                except Exception as citation_error:
                    logger.warning(f"Error extracting citations for cached summary: {citation_error}")
                    citations = None
                
                return NotebookQueryResponse(
                    answer=cached_summary["answer"],
                    mode=cached_summary["mode"],
                    context_used=cached_summary["context_used"],
                    citations=citations,
                    source_documents=None,
                    chat_context_used=False
                )
            
            # Generate new summary if no valid cache exists
            logger.info(f"Generating new summary for notebook {notebook_id} (documents changed)")
            
            # Get existing RAG instance
            rag = await get_lightrag_instance(notebook_id)
            
            # Create summary query
            summary_question = ("Write down a comprehensive summary of all the documents provided in a single paragraph. "
                              "Mention what the documents are about, the main topics they cover, key themes, "
                              "important findings or insights, and the overall scope of the content. "
                              "Focus on providing an overview that helps understand the nature and breadth of the knowledge base.")
            
            # Create query parameters optimized for summary generation
            query_param = QueryParam(
                mode="hybrid",  # Use hybrid mode for comprehensive coverage
                response_type="Single Paragraph",  # Request single paragraph format
                top_k=100,  # Use higher top_k to get broader coverage of documents
            )
            
            # Perform summary query
            result = await rag.aquery(summary_question, param=query_param)
            
            # Extract citation information for all completed documents
            citations = []
            try:
                for doc in notebook_documents:
                    citation = {
                        "filename": doc["filename"],
                        "file_path": doc.get("file_path", f"documents/{doc['filename']}"),
                        "document_id": doc["id"],
                        "title": doc["filename"].replace('_', ' ').replace('.txt', '').replace('.pdf', '').replace('.md', '').title()
                    }
                    citations.append(citation)
                
            except Exception as citation_error:
                logger.warning(f"Error extracting citations for summary: {citation_error}")
                citations = None
            
            # Cache the new summary
            summary_cache = {
                "answer": result,
                "mode": "hybrid",
                "context_used": True,
                "generated_at": datetime.now().isoformat()
            }
            
            # Update notebook with cached summary and fingerprint
            lightrag_notebooks_db[notebook_id]["summary_cache"] = summary_cache
            lightrag_notebooks_db[notebook_id]["docs_fingerprint"] = current_docs_fingerprint
            
            # Save to disk
            save_notebooks_db()
            
            logger.info(f"Generated and cached new summary for notebook {notebook_id} with {len(notebook_documents)} documents")
            
            return NotebookQueryResponse(
                answer=result,
                mode="hybrid",
                context_used=True,
                citations=citations,
                source_documents=None,
                chat_context_used=False
            )
            
        except Exception as e:
            logger.error(f"Error generating summary for notebook {notebook_id}: {e}")
            raise HTTPException(status_code=500, detail=f"Error generating summary: {str(e)}")

    @app.get("/notebooks/{notebook_id}/graph")
    async def get_notebook_graph_data(notebook_id: str):
        """Get graph visualization data for a notebook"""
        validate_notebook_exists(notebook_id)
        
        try:
            # Path to the GraphML file created by LightRAG
            working_dir = LIGHTRAG_STORAGE_PATH / notebook_id
            graphml_file = working_dir / "graph_chunk_entity_relation.graphml"
            
            if not graphml_file.exists():
                return {
                    "nodes": [],
                    "edges": [],
                    "message": "No graph data available yet. Upload documents and query the notebook to generate the knowledge graph."
                }
            
            # Parse the GraphML file
            tree = ET.parse(graphml_file)
            root = tree.getroot()
            
            # GraphML namespace
            ns = {'graphml': 'http://graphml.graphdrawing.org/xmlns'}
            
            nodes = []
            edges = []
            
            # Extract nodes
            for node in root.findall('.//graphml:node', ns):
                node_id = node.get('id')
                node_data = {'id': node_id, 'type': 'entity', 'properties': {}}
                
                # Extract node attributes/data
                for data in node.findall('graphml:data', ns):
                    key = data.get('key')
                    value = data.text or ''
                    
                    # Map common GraphML keys to readable labels
                    if key == 'd0':  # Usually entity name
                        node_data['label'] = value
                        node_data['properties']['name'] = value
                    elif key == 'd1':  # Usually entity type
                        node_data['type'] = value
                        node_data['properties']['entity_type'] = value
                    elif key == 'd2':  # Usually description
                        node_data['properties']['description'] = value
                    else:
                        node_data['properties'][key] = value
                
                # Use node_id as label if no label found
                if 'label' not in node_data:
                    node_data['label'] = node_id
                
                nodes.append(node_data)
            
            # Extract edges
            for edge in root.findall('.//graphml:edge', ns):
                source = edge.get('source')
                target = edge.get('target')
                edge_data = {
                    'source': source,
                    'target': target,
                    'relationship': 'related_to',
                    'properties': {}
                }
                
                # Extract edge attributes/data
                for data in edge.findall('graphml:data', ns):
                    key = data.get('key')
                    value = data.text or ''
                    
                    # Map common GraphML keys for edges
                    if key == 'd3':  # Usually relationship type
                        edge_data['relationship'] = value
                        edge_data['properties']['relation_type'] = value
                    elif key == 'd4':  # Usually weight or strength
                        try:
                            edge_data['properties']['weight'] = float(value)
                        except ValueError:
                            edge_data['properties']['weight_str'] = value
                    elif key == 'd5':  # Usually description
                        edge_data['properties']['description'] = value
                    else:
                        edge_data['properties'][key] = value
                
                edges.append(edge_data)
            
            logger.info(f"Loaded graph data for notebook {notebook_id}: {len(nodes)} nodes, {len(edges)} edges")
            
            return {
                "nodes": nodes,
                "edges": edges,
                "stats": {
                    "node_count": len(nodes),
                    "edge_count": len(edges),
                    "file_path": str(graphml_file)
                }
            }
            
        except ET.ParseError as e:
            logger.error(f"Error parsing GraphML file for notebook {notebook_id}: {e}")
            return {
                "nodes": [],
                "edges": [],
                "error": f"Error parsing graph data: {str(e)}"
            }
        except Exception as e:
            logger.error(f"Error getting graph data for notebook {notebook_id}: {e}")
            raise HTTPException(status_code=500, detail=f"Error getting graph data: {str(e)}")

    @app.get("/notebooks/{notebook_id}/graph/html")
    async def get_notebook_graph_html(notebook_id: str):
        """Generate interactive HTML graph visualization using pyvis and networkx"""
        validate_notebook_exists(notebook_id)
        
        try:
            # Path to the GraphML file created by LightRAG
            working_dir = LIGHTRAG_STORAGE_PATH / notebook_id
            graphml_file = working_dir / "graph_chunk_entity_relation.graphml"
            
            if not graphml_file.exists():
                # Return a simple HTML page indicating no data
                html_content = """
                <!DOCTYPE html>
                <html>
                <head>
                    <title>Knowledge Graph - No Data</title>
                    <style>
                        body { font-family: Arial, sans-serif; text-align: center; padding: 50px; }
                        .message { color: #666; font-size: 18px; }
                    </style>
                </head>
                <body>
                    <h2>Knowledge Graph</h2>
                    <p class="message">No graph data available yet.</p>
                    <p class="message">Upload documents and query the notebook to generate the knowledge graph.</p>
                </body>
                </html>
                """
                return HTMLResponse(content=html_content)
            
            # Install required packages if not available
            try:
                import networkx as nx
                from pyvis.network import Network
            except ImportError:
                # Try to install packages
                import subprocess
                import sys
                
                logger.info("Installing required packages for graph visualization...")
                try:
                    subprocess.check_call([sys.executable, "-m", "pip", "install", "pyvis", "networkx"])
                    import networkx as nx
                    from pyvis.network import Network
                    logger.info("Successfully installed pyvis and networkx")
                except Exception as install_error:
                    logger.error(f"Failed to install required packages: {install_error}")
                    raise HTTPException(status_code=500, detail="Required packages (pyvis, networkx) not available")
            
            import random
            import tempfile
            
            # Load the GraphML file
            logger.info(f"Loading GraphML file: {graphml_file}")
            G = nx.read_graphml(str(graphml_file))
            
            # Create a Pyvis network with responsive design
            net = Network(
                height="100vh",
                width="100%",
                bgcolor="#ffffff",
                font_color="#333333",
                notebook=False
            )
            
            # Configure physics for better layout
            net.set_options("""
            var options = {
              "physics": {
                "enabled": true,
                "stabilization": {"iterations": 100},
                "barnesHut": {
                  "gravitationalConstant": -8000,
                  "centralGravity": 0.3,
                  "springLength": 95,
                  "springConstant": 0.04,
                  "damping": 0.09
                }
              },
              "nodes": {
                "font": {"size": 12},
                "scaling": {
                  "min": 10,
                  "max": 30
                }
              },
              "edges": {
                "font": {"size": 10},
                "scaling": {
                  "min": 1,
                  "max": 3
                }
              }
            }
            """)
            
            # Convert NetworkX graph to Pyvis network
            net.from_nx(G)
            
            # Define colors for different node types
            node_type_colors = {
                'person': '#FF6B6B',      # Red
                'organization': '#4ECDC4', # Teal
                'location': '#45B7D1',    # Blue
                'concept': '#96CEB4',     # Green
                'event': '#FFEAA7',      # Yellow
                'entity': '#DDA0DD',     # Plum
                'default': '#95A5A6'     # Gray
            }
            
            # Enhance nodes with colors, titles, and better styling
            for node in net.nodes:
                # Determine node type from the data
                node_type = 'entity'  # default
                if 'entity_type' in node:
                    node_type = str(node['entity_type']).lower()
                elif 'type' in node:
                    node_type = str(node['type']).lower()
                
                # Set color based on type
                node["color"] = node_type_colors.get(node_type, node_type_colors['default'])
                
                # Add hover title with description
                title_parts = [f"ID: {node.get('id', 'Unknown')}"]
                if 'label' in node and node['label']:
                    title_parts.append(f"Label: {node['label']}")
                if node_type:
                    title_parts.append(f"Type: {node_type.title()}")
                if 'description' in node and node['description']:
                    desc = str(node['description'])[:200] + "..." if len(str(node['description'])) > 200 else str(node['description'])
                    title_parts.append(f"Description: {desc}")
                
                node["title"] = "\\n".join(title_parts)
                
                # Set node size based on connections (degree)
                if hasattr(G, 'degree'):
                    degree = G.degree(node['id']) if node['id'] in G else 1
                    node["size"] = min(10 + degree * 2, 30)  # Size between 10-30
                
                # Clean up label for display
                if 'label' in node and node['label']:
                    # Truncate long labels
                    label = str(node['label'])
                    node["label"] = label[:20] + "..." if len(label) > 20 else label
                else:
                    # Use ID as label if no label exists
                    node_id = str(node.get('id', ''))
                    node["label"] = node_id[:20] + "..." if len(node_id) > 20 else node_id
            
            # Enhance edges with titles and styling
            for edge in net.edges:
                title_parts = []
                
                # Add relationship type
                if 'relationship' in edge and edge['relationship']:
                    title_parts.append(f"Relationship: {edge['relationship']}")
                elif 'relation_type' in edge and edge['relation_type']:
                    title_parts.append(f"Relationship: {edge['relation_type']}")
                
                # Add weight if available
                if 'weight' in edge and edge['weight']:
                    try:
                        weight = float(edge['weight'])
                        title_parts.append(f"Weight: {weight:.2f}")
                        # Set edge width based on weight
                        edge["width"] = min(max(1, weight * 2), 5)
                    except (ValueError, TypeError):
                        pass
                
                # Add description if available
                if 'description' in edge and edge['description']:
                    desc = str(edge['description'])[:100] + "..." if len(str(edge['description'])) > 100 else str(edge['description'])
                    title_parts.append(f"Description: {desc}")
                
                if title_parts:
                    edge["title"] = "\\n".join(title_parts)
                
                # Style edges
                edge["color"] = {"color": "#848484", "highlight": "#333333"}
            
            # Generate HTML
            with tempfile.NamedTemporaryFile(mode='w', suffix='.html', delete=False) as tmp_file:
                net.save_graph(tmp_file.name)
                tmp_file.flush()
                
                # Read the generated HTML
                with open(tmp_file.name, 'r', encoding='utf-8') as f:
                    html_content = f.read()
                
                # Clean up temp file
                import os
                os.unlink(tmp_file.name)
            
            # Enhance the HTML with custom styling and dark mode support
            enhanced_html = html_content.replace(
                '<head>',
                '''<head>
                <style>
                    body { 
                        margin: 0; 
                        padding: 0; 
                        font-family: Arial, sans-serif;
                        background: #f8f9fa;
                    }
                    
                    @media (prefers-color-scheme: dark) {
                        body { background: #1a1a1a; }
                    }
                    
                    .graph-container {
                        position: relative;
                        width: 100%;
                        height: 100vh;
                    }
                    
                    .graph-info {
                        position: absolute;
                        top: 10px;
                        left: 10px;
                        background: rgba(255, 255, 255, 0.9);
                        padding: 10px;
                        border-radius: 5px;
                        font-size: 12px;
                        z-index: 1000;
                        box-shadow: 0 2px 5px rgba(0,0,0,0.1);
                    }
                    
                    @media (prefers-color-scheme: dark) {
                        .graph-info { 
                            background: rgba(30, 30, 30, 0.9); 
                            color: white;
                        }
                    }
                    
                    .legend {
                        position: absolute;
                        top: 10px;
                        right: 10px;
                        background: rgba(255, 255, 255, 0.9);
                        padding: 10px;
                        border-radius: 5px;
                        font-size: 11px;
                        z-index: 1000;
                        box-shadow: 0 2px 5px rgba(0,0,0,0.1);
                        max-width: 200px;
                    }
                    
                    @media (prefers-color-scheme: dark) {
                        .legend { 
                            background: rgba(30, 30, 30, 0.9); 
                            color: white;
                        }
                    }
                    
                    .legend-item {
                        display: flex;
                        align-items: center;
                        margin: 2px 0;
                    }
                    
                    .legend-color {
                        width: 12px;
                        height: 12px;
                        border-radius: 50%;
                        margin-right: 5px;
                    }
                </style>'''
            )
            
            # Add info overlay and legend
            graph_stats = f"Nodes: {len(net.nodes)} | Edges: {len(net.edges)}"
            
            legend_html = '''
            <div class="legend">
                <strong>Node Types:</strong>
                <div class="legend-item">
                    <div class="legend-color" style="background: #FF6B6B;"></div>
                    <span>Person</span>
                </div>
                <div class="legend-item">
                    <div class="legend-color" style="background: #4ECDC4;"></div>
                    <span>Organization</span>
                </div>
                <div class="legend-item">
                    <div class="legend-color" style="background: #45B7D1;"></div>
                    <span>Location</span>
                </div>
                <div class="legend-item">
                    <div class="legend-color" style="background: #96CEB4;"></div>
                    <span>Concept</span>
                </div>
                <div class="legend-item">
                    <div class="legend-color" style="background: #FFEAA7;"></div>
                    <span>Event</span>
                </div>
                <div class="legend-item">
                    <div class="legend-color" style="background: #DDA0DD;"></div>
                    <span>Entity</span>
                </div>
            </div>
            '''
            
            enhanced_html = enhanced_html.replace(
                '<body>',
                f'''<body>
                <div class="graph-info">{graph_stats}</div>
                {legend_html}'''
            )
            
            logger.info(f"Generated interactive graph HTML for notebook {notebook_id}: {len(net.nodes)} nodes, {len(net.edges)} edges")
            
            return HTMLResponse(content=enhanced_html)
            
        except Exception as e:
            logger.error(f"Error generating graph HTML for notebook {notebook_id}: {e}")
            logger.error(f"Full traceback: {traceback.format_exc()}")
            
            # Return error HTML page
            error_html = f"""
            <!DOCTYPE html>
            <html>
            <head>
                <title>Knowledge Graph - Error</title>
                <style>
                    body {{ font-family: Arial, sans-serif; text-align: center; padding: 50px; }}
                    .error {{ color: #d32f2f; font-size: 18px; }}
                    .details {{ color: #666; font-size: 14px; margin-top: 10px; }}
                </style>
            </head>
            <body>
                <h2>Knowledge Graph</h2>
                <p class="error">Error generating graph visualization</p>
                <p class="details">{str(e)}</p>
            </body>
            </html>
            """
            return HTMLResponse(content=error_html)

    @app.get("/notebooks/{notebook_id}/debug")
    async def debug_notebook_documents(notebook_id: str):
        """Debug endpoint to check document processing status and LightRAG state"""
        validate_notebook_exists(notebook_id)
        
        try:
            # Get all documents for this notebook
            notebook_documents = [
                doc for doc in lightrag_documents_db.values() 
                if doc["notebook_id"] == notebook_id
            ]
            
            # Get LightRAG instance info
            rag_info = {"exists": False, "working_dir": None}
            if notebook_id in lightrag_instances:
                rag = lightrag_instances[notebook_id]
                working_dir = LIGHTRAG_STORAGE_PATH / notebook_id
                rag_info = {
                    "exists": True,
                    "working_dir": str(working_dir),
                    "directory_exists": working_dir.exists(),
                    "files": list(working_dir.glob("*")) if working_dir.exists() else []
                }
            
            return {
                "notebook_id": notebook_id,
                "documents_count": len(notebook_documents),
                "documents": [
                    {
                        "id": doc["id"],
                        "filename": doc["filename"],
                        "status": doc["status"],
                        "lightrag_id": doc.get("lightrag_id"),
                        "error": doc.get("error"),
                        "uploaded_at": doc["uploaded_at"].isoformat() if isinstance(doc["uploaded_at"], datetime) else doc["uploaded_at"]
                    }
                    for doc in notebook_documents
                ],
                "lightrag_info": rag_info
            }
            
        except Exception as e:
            logger.error(f"Error in debug endpoint: {e}")
            raise HTTPException(status_code=500, detail=f"Debug error: {str(e)}")

    # Chat History Endpoints
    @app.get("/notebooks/{notebook_id}/chat/history", response_model=ChatHistoryResponse)
    async def get_chat_history(notebook_id: str, limit: int = Query(50, description="Maximum number of messages to return")):
        """Get chat history for a notebook"""
        validate_notebook_exists(notebook_id)
        
        messages = chat_history_db.get(notebook_id, [])
        
        # Limit messages and convert to ChatMessage objects
        limited_messages = messages[-limit:] if limit > 0 else messages
        chat_messages = [ChatMessage(**msg) for msg in limited_messages]
        
        return ChatHistoryResponse(
            notebook_id=notebook_id,
            messages=chat_messages,
            total_messages=len(messages)
        )

    @app.delete("/notebooks/{notebook_id}/chat/history")
    async def clear_chat_history(notebook_id: str):
        """Clear chat history for a notebook"""
        validate_notebook_exists(notebook_id)
        
        if notebook_id in chat_history_db:
            del chat_history_db[notebook_id]
            save_chat_history_db()
        
        return {"message": "Chat history cleared successfully"}

    # Enhanced Query with Chat History
    @app.post("/notebooks/{notebook_id}/chat", response_model=NotebookQueryResponse)
    async def chat_with_notebook(notebook_id: str, query: NotebookQueryRequest):
        """Chat with a notebook using conversation history"""
        validate_notebook_exists(notebook_id)
        
        try:
            # Initialize chat history if it doesn't exist
            if notebook_id not in chat_history_db:
                chat_history_db[notebook_id] = []
            
            # Add user message to history
            user_message = {
                "role": "user",
                "content": query.question,
                "timestamp": datetime.now()
            }
            chat_history_db[notebook_id].append(user_message)
            
            # Get LightRAG instance
            rag = await get_lightrag_instance(notebook_id)
            
            # Build context from chat history if enabled
            chat_context = ""
            if query.use_chat_history and len(chat_history_db[notebook_id]) > 1:
                recent_messages = chat_history_db[notebook_id][-10:]  # Last 10 messages
                chat_context = "Previous conversation context:\n"
                for msg in recent_messages[:-1]:  # Exclude the current message
                    chat_context += f"{msg['role'].title()}: {msg['content']}\n"
                chat_context += "\nCurrent question: "
            
            # Enhance question with chat context
            enhanced_question = chat_context + query.question if chat_context else query.question
            
            # Map mode string to QueryParam
            mode_mapping = {
                "naive": "naive",
                "local": "local", 
                "global": "global",
                "hybrid": "hybrid",
                "mix": "mix"
            }
            adjusted_mode = mode_mapping.get(query.mode, "hybrid")
            
            # Execute query
            query_param = QueryParam(
                mode=adjusted_mode,
                response_type=query.response_type,
                top_k=query.top_k
            )
            
            result = await rag.aquery(enhanced_question, param=query_param)
            
            # Extract citations (basic implementation)
            citations = []
            try:
                # This is a basic implementation - can be enhanced based on LightRAG output format
                citations = [{"source": "Document Context", "relevance": "High"}]
            except Exception as citation_error:
                logger.warning(f"Error extracting citations: {citation_error}")
            
            # Add assistant message to history
            assistant_message = {
                "role": "assistant",
                "content": str(result),
                "timestamp": datetime.now(),
                "citations": citations
            }
            chat_history_db[notebook_id].append(assistant_message)
            
            # Save chat history
            save_chat_history_db()
            
            return NotebookQueryResponse(
                answer=str(result),
                mode=adjusted_mode,
                context_used=True,
                citations=citations,
                source_documents=None,  # Will be enhanced later
                chat_context_used=bool(chat_context)
            )
            
        except Exception as e:
            logger.error(f"Error in chat query for notebook {notebook_id}: {e}")
            logger.error(f"Full error traceback: {traceback.format_exc()}")
            raise HTTPException(status_code=500, detail=f"Error processing chat query: {str(e)}")

    # Document Summary Endpoints
    @app.post("/notebooks/{notebook_id}/summary/detailed", response_model=NotebookQueryResponse)
    async def generate_detailed_summary(notebook_id: str, request: DocumentSummaryRequest):
        """Generate a detailed summary with document-level insights"""
        validate_notebook_exists(notebook_id)
        
        try:
            # Get document list
            notebook_documents = [
                doc for doc in lightrag_documents_db.values() 
                if doc["notebook_id"] == notebook_id and doc["status"] == "completed"
            ]
            
            if not notebook_documents:
                return NotebookQueryResponse(
                    answer="No completed documents found. Please upload and process documents first.",
                    mode="summary",
                    context_used=False,
                    citations=None,
                    source_documents=None,
                    chat_context_used=False
                )
            
            # Get LightRAG instance
            rag = await get_lightrag_instance(notebook_id)
            
            # Create detailed summary prompt based on length preference
            length_prompts = {
                "short": "Write a concise summary (2-3 sentences) of the key points from all documents.",
                "medium": "Write a comprehensive summary (1-2 paragraphs) covering the main themes and insights from all documents.",
                "long": "Write a detailed analysis (3-4 paragraphs) covering key themes, insights, contradictions, and conclusions from all documents."
            }
            
            summary_prompt = length_prompts.get(request.max_length, length_prompts["medium"])
            
            if request.include_details:
                summary_prompt += f"\n\nInclude insights from these {len(notebook_documents)} documents: " + \
                                ", ".join([doc["filename"] for doc in notebook_documents])
            
            # Execute summary query
            query_param = QueryParam(mode="hybrid", response_type="Multiple Paragraphs", top_k=100)
            result = await rag.aquery(summary_prompt, param=query_param)
            
            # Build source documents list
            source_docs = []
            if request.include_details:
                for doc in notebook_documents:
                    source_docs.append({
                        "filename": doc["filename"],
                        "upload_date": doc["uploaded_at"].isoformat() if isinstance(doc["uploaded_at"], datetime) else doc["uploaded_at"],
                        "status": doc["status"]
                    })
            
            return NotebookQueryResponse(
                answer=str(result),
                mode="hybrid",
                context_used=True,
                citations=[{"source": doc["filename"], "type": "document"} for doc in notebook_documents],
                source_documents=source_docs,
                chat_context_used=False
            )
            
        except Exception as e:
            logger.error(f"Error generating detailed summary for notebook {notebook_id}: {e}")
            raise HTTPException(status_code=500, detail=f"Error generating summary: {str(e)}")

    # Query Templates Endpoint
    @app.get("/query-templates", response_model=List[QueryTemplate])
    async def get_query_templates():
        """Get pre-built query templates for common use cases"""
        templates = [
            QueryTemplate(
                id="summarize_all",
                name="Summarize All Documents",
                description="Get a comprehensive overview of all uploaded documents",
                question_template="Provide a comprehensive summary of all the documents, highlighting the main themes and key insights.",
                category="Analysis",
                use_case="When you want a high-level overview of your knowledge base"
            ),
            QueryTemplate(
                id="find_contradictions",
                name="Find Contradictions",
                description="Identify conflicting information across documents",
                question_template="Are there any contradictions or conflicting viewpoints between the different documents? Please explain any differences you find.",
                category="Analysis",
                use_case="When comparing different sources or viewpoints"
            ),
            QueryTemplate(
                id="extract_key_facts",
                name="Extract Key Facts",
                description="Pull out the most important facts and data points",
                question_template="What are the most important facts, statistics, and data points mentioned across all documents?",
                category="Research",
                use_case="When you need specific factual information"
            ),
            QueryTemplate(
                id="timeline_analysis",
                name="Timeline Analysis",
                description="Understand chronological order and progression",
                question_template="Can you create a timeline of events or developments mentioned in the documents?",
                category="Research",
                use_case="When tracking historical progression or project timelines"
            ),
            QueryTemplate(
                id="compare_approaches",
                name="Compare Approaches",
                description="Compare different methods or strategies",
                question_template="What are the different approaches or methodologies discussed in the documents? Compare their advantages and disadvantages.",
                category="Comparison",
                use_case="When evaluating different options or strategies"
            ),
            QueryTemplate(
                id="action_items",
                name="Extract Action Items",
                description="Find actionable recommendations and next steps",
                question_template="What are the key recommendations, action items, or next steps suggested in the documents?",
                category="Planning",
                use_case="When planning follow-up activities"
            ),
            QueryTemplate(
                id="expert_opinions",
                name="Expert Opinions",
                description="Identify expert viewpoints and authoritative sources",
                question_template="What do experts or authoritative sources say about this topic? Include any credentials or authority mentioned.",
                category="Research",
                use_case="When you need authoritative perspectives"
            ),
            QueryTemplate(
                id="gaps_analysis",
                name="Knowledge Gaps",
                description="Identify missing information or areas that need more research",
                question_template="What topics or questions are mentioned but not fully explained? What knowledge gaps exist that might need additional research?",
                category="Analysis",
                use_case="When planning additional research"
            ),
            QueryTemplate(
                id="practical_applications",
                name="Practical Applications",
                description="Find real-world applications and use cases",
                question_template="What are the practical applications or real-world use cases mentioned in the documents? Include any examples or case studies.",
                category="Application",
                use_case="When looking for implementation ideas"
            ),
            QueryTemplate(
                id="risk_analysis",
                name="Risk Analysis",
                description="Identify potential risks, challenges, or limitations",
                question_template="What risks, challenges, limitations, or potential problems are mentioned in the documents?",
                category="Analysis",
                use_case="When assessing potential issues or preparing for challenges"
            )
        ]
        
        return templates

    @app.post("/notebooks/{notebook_id}/query/template/{template_id}", response_model=NotebookQueryResponse)
    async def query_with_template(notebook_id: str, template_id: str, custom_params: Optional[Dict[str, str]] = None):
        """Execute a query using a pre-built template"""
        validate_notebook_exists(notebook_id)
        
        # Get templates
        templates = await get_query_templates()
        template = next((t for t in templates if t.id == template_id), None)
        
        if not template:
            raise HTTPException(status_code=404, detail="Template not found")
        
        # Build question from template
        question = template.question_template
        
        # Replace any custom parameters if provided
        if custom_params:
            for key, value in custom_params.items():
                question = question.replace(f"{{{key}}}", value)
        
        # Create query request
        query_request = NotebookQueryRequest(
            question=question,
            mode="hybrid",
            response_type="Multiple Paragraphs",
            top_k=60,
            llm_provider=None,  # Use notebook's default provider
            use_chat_history=False  # Templates don't use chat history by default
        )
        
        # Execute using the chat endpoint for consistency
        return await chat_with_notebook(notebook_id, query_request)

# Document Text Extraction Endpoint
@app.post("/extract-text")
async def extract_document_text(file: UploadFile = File(...)):
    """
    Extract text from a document file without creating a notebook.
    Supports: PDF, TXT, MD, RTF, DOC, DOCX, XLS, XLSX, PPT, PPTX, HTML, HTM, XML, CSV, JSON
    """
    if not file.filename:
        raise HTTPException(status_code=400, detail="No filename provided")
    
    # Check file size (50MB limit)
    max_file_size = 50 * 1024 * 1024  # 50MB
    file_content = await file.read()
    
    if len(file_content) > max_file_size:
        raise HTTPException(
            status_code=400, 
            detail=f"File too large. Maximum size is 50MB, received {len(file_content) / (1024*1024):.1f}MB"
        )
    
    if not file_content:
        raise HTTPException(status_code=400, detail="Empty file")
    
    try:
        # Extract text using the existing function
        extracted_text = await extract_text_from_file(file.filename, file_content)
        
        # Get file statistics
        file_extension = file.filename.lower().split('.')[-1] if '.' in file.filename else 'unknown'
        text_length = len(extracted_text)
        word_count = len(extracted_text.split()) if extracted_text else 0
        
        return {
            "status": "success",
            "filename": file.filename,
            "file_type": file_extension.upper(),
            "file_size_bytes": len(file_content),
            "file_size_mb": round(len(file_content) / (1024*1024), 2),
            "extracted_text": extracted_text,
            "text_stats": {
                "character_count": text_length,
                "word_count": word_count,
                "line_count": extracted_text.count('\n') + 1 if extracted_text else 0
            }
        }
        
    except HTTPException:
        # Re-raise HTTP exceptions (like unsupported file types)
        raise
    except Exception as e:
        logger.error(f"Error extracting text from {file.filename}: {e}")
        logger.error(traceback.format_exc())
        raise HTTPException(
            status_code=500, 
            detail=f"Error processing document: {str(e)}"
        )

# Audio transcription endpoint
@app.post("/transcribe")
async def transcribe_audio(
    file: UploadFile = File(...),
    language: Optional[str] = Form(None),
    beam_size: int = Form(5),
    initial_prompt: Optional[str] = Form(None)
):
    """Transcribe an audio file using faster-whisper (CPU mode)"""
    # Validate file extension
    supported_formats = ['mp3', 'wav', 'flac', 'm4a', 'ogg', 'opus']
    
    if not file.filename:
        raise HTTPException(status_code=400, detail="No filename provided")
    
    file_extension = file.filename.lower().split('.')[-1]
    
    if file_extension not in supported_formats:
        raise HTTPException(
            status_code=400, 
            detail=f"Unsupported audio format: {file_extension}. Supported formats: {', '.join(supported_formats)}"
        )
    
    # Read file content
    try:
        content = await file.read()
        if not content:
            raise HTTPException(status_code=400, detail="Empty audio file")
    except Exception as e:
        logger.error(f"Error reading audio file: {e}")
        raise HTTPException(status_code=500, detail=f"Error reading audio file: {str(e)}")
    
    # Get Speech2Text instance
    try:
        s2t = get_speech2text()
    except Exception as e:
        logger.error(f"Error initializing Speech2Text: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to initialize Speech2Text: {str(e)}")
    
    # Transcribe the audio
    try:
        result = s2t.transcribe_bytes(
            content,
            language=language,
            beam_size=beam_size,
            initial_prompt=initial_prompt
        )
        
        return {
            "status": "success",
            "filename": file.filename,
            "transcription": result
        }
    except Exception as e:
        logger.error(f"Error transcribing audio: {e}")
        logger.error(traceback.format_exc())
        raise HTTPException(status_code=500, detail=f"Error transcribing audio: {str(e)}")

# Text-to-Speech endpoints
@app.post("/synthesize")
async def synthesize_text(request: TTSRequest):
    """
    Synthesize text to speech and return audio data.
    Supports multiple TTS engines including Kokoro for high-quality neural TTS.
    """
    try:
        logger.info(f"TTS request: text='{request.text[:50]}...', engine={request.engine}, voice={request.voice}, speed={request.speed}")
        
        # Get Text2Speech instance
        t2s = get_text2speech()
        
        # If specific engine requested, create new instance with those settings
        if request.engine and request.engine != t2s.engine:
            logger.info(f"Creating new TTS instance with engine: {request.engine}")
            t2s = Text2Speech(
                engine=request.engine,
                language=request.language or "en",
                slow=request.slow or False,
                voice=request.voice or "af_sarah",
                speed=request.speed or 1.0
            )
        
        # Generate speech
        audio_bytes = t2s.synthesize_to_bytes(request.text)
        
        # Determine content type based on engine
        if request.engine in ["kokoro", "kokoro-onnx"]:
            content_type = "audio/wav"
        else:
            content_type = "audio/mpeg"
        
        return Response(
            content=audio_bytes,
            media_type=content_type,
            headers={
                "Content-Disposition": "attachment; filename=speech.wav" if request.engine in ["kokoro", "kokoro-onnx"] else "attachment; filename=speech.mp3"
            }
        )
        
    except Exception as e:
        logger.error(f"Error in TTS synthesis: {e}")
        logger.error(traceback.format_exc())
        raise HTTPException(status_code=500, detail=f"TTS synthesis failed: {str(e)}")

@app.post("/synthesize/file")
async def synthesize_text_to_file(
    text: str = Form(...),
    language: Optional[str] = Form("en"),
    engine: Optional[str] = Form(None),
    slow: Optional[bool] = Form(False),
    voice: Optional[str] = Form("af_sarah"),
    speed: Optional[float] = Form(1.0),
    filename: Optional[str] = Form("speech.mp3")
):
    """
    Synthesize text to speech and return as downloadable file.
    Supports multiple TTS engines including Kokoro for high-quality neural TTS.
    """
    try:
        logger.info(f"TTS file request: text='{text[:50]}...', engine={engine}, voice={voice}, speed={speed}")
        
        # Get Text2Speech instance
        t2s = get_text2speech()
        
        # If a specific engine is requested and different from current
        if engine and engine != t2s.engine:
            logger.info(f"Creating new TTS instance with engine: {engine}")
            t2s = Text2Speech(
                engine=engine,
                language=language or "en",
                slow=slow or False,
                voice=voice or "af_sarah",
                speed=speed or 1.0
            )
        
        # Determine file extension based on engine
        if engine in ["kokoro", "kokoro-onnx"]:
            file_ext = ".wav"
            content_type = "audio/wav"
        else:
            file_ext = ".mp3"
            content_type = "audio/mpeg"
        
        # Ensure filename has correct extension
        if not filename:
            filename = "speech" + file_ext
        elif not filename.endswith(file_ext):
            filename = os.path.splitext(filename)[0] + file_ext
        
        # Create temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix=file_ext) as temp_file:
            temp_path = temp_file.name
        
        try:
            # Generate speech to file
            output_path = t2s.synthesize_to_file(text, temp_path)
            
            # Read the file
            with open(output_path, 'rb') as f:
                audio_data = f.read()
            
            return Response(
                content=audio_data,
                media_type=content_type,
                headers={
                    "Content-Disposition": f"attachment; filename={filename}",
                    "Content-Length": str(len(audio_data))
                }
            )
            
        finally:
            # Clean up temp file
            try:
                os.unlink(temp_path)
            except:
                pass
        
    except Exception as e:
        logger.error(f"Error in TTS file synthesis: {e}")
        logger.error(traceback.format_exc())
        raise HTTPException(status_code=500, detail=f"TTS file synthesis failed: {str(e)}")

@app.get("/tts/languages")
async def get_tts_languages():
    """Get available languages for text-to-speech"""
    try:
        t2s = get_text2speech()
        languages = t2s.get_available_languages()
        
        return {
            "engine": t2s.engine,
            "current_language": t2s.language,
            "available_languages": languages
        }
    except Exception as e:
        logger.error(f"Error getting TTS languages: {e}")
        raise HTTPException(status_code=500, detail=f"Error getting TTS languages: {str(e)}")

@app.get("/tts/status")
async def get_tts_status():
    """Get current TTS engine status and configuration"""
    try:
        t2s = get_text2speech()
        
        return {
            "engine": t2s.engine,
            "language": t2s.language,
            "slow": t2s.slow,
            "available_engines": []
        }
    except Exception as e:
        logger.error(f"Error getting TTS status: {e}")
        raise HTTPException(status_code=500, detail=f"Error getting TTS status: {str(e)}")

@app.get("/tts/voices")
async def get_tts_voices():
    """Get available voices for TTS engines"""
    try:
        voices = {
            "kokoro_voices": {
                "af_sarah": "American Female - Sarah (warm, friendly)",
                "af_nicole": "American Female - Nicole (professional)",
                "af_sky": "American Female - Sky (energetic)",
                "am_adam": "American Male - Adam (deep, authoritative)",
                "am_michael": "American Male - Michael (casual)",
                "bf_emma": "British Female - Emma (elegant)",
                "bf_isabella": "British Female - Isabella (sophisticated)",
                "bm_george": "British Male - George (distinguished)",
                "bm_lewis": "British Male - Lewis (modern)"
            },
            "pyttsx3_voices": "System dependent - use /tts/status to see available voices",
            "gtts_languages": [
                "en", "es", "fr", "de", "it", "pt", "ru", "ja", "ko", "zh", "hi", "ar"
            ]
        }
        
        # Try to get actual pyttsx3 voices if available
        try:
            import pyttsx3
            engine = pyttsx3.init()
            system_voices = engine.getProperty('voices')
            if system_voices:
                voices["pyttsx3_voices"] = [
                    {
                        "id": voice.id,
                        "name": voice.name,
                        "languages": getattr(voice, 'languages', []),
                        "gender": getattr(voice, 'gender', 'unknown')
                    }
                    for voice in system_voices
                ]
            engine.stop()
        except:
            pass
        
        return voices
        
    except Exception as e:
        logger.error(f"Error getting TTS voices: {e}")
        raise HTTPException(status_code=500, detail=f"Error getting TTS voices: {str(e)}")

# Handle graceful shutdown
def handle_exit(signum, frame):
    logger.info(f"Received signal {signum}, shutting down gracefully")
    sys.exit(0)

# Register signal handlers
signal.signal(signal.SIGINT, handle_exit)
signal.signal(signal.SIGTERM, handle_exit)

if __name__ == "__main__":
    import uvicorn
    logger.info(f"Starting server on {HOST}:{PORT}")
    
    # Start the server with reload=False to prevent duplicate processes
    uvicorn.run(
        "main:app",
        host=HOST,
        port=PORT,
        log_level="info",
        reload=False  # Change this to false to prevent multiple processes
    )
